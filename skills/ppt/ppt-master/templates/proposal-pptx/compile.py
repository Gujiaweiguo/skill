#!/usr/bin/env python3
"""
方案介绍PPT模板编译器 (Proposal PPT Template Compiler)

基于参考产品PPT + 客户定制分析页 → 生成完整方案介绍PPT

架构:
  参考 PPT (62页, 产品截图/案例) = 固定基底
  YAML 内容包 (客户需求/覆盖度/场景/报价/计划) = 定制数据
  compile.py = python-pptx 编译器
  输出 = 基底 + 5张定制页 (插入封面之后)

定制页型:
  1. requirement-understanding  需求理解 (分类卡片 + 核心洞察)
  2. coverage-analysis          功能覆盖度 (三大指标卡)
  3. core-scenario              核心场景 (双行闭环流)
  4. pricing-comparison         报价对比 (双方案卡片)
  5. implementation-plan        实施计划 (阶段时间线)

Usage:
  uv run compile.py <content-package.yaml>
  uv run compile.py <content-package.yaml> --base ref.pptx --output out.pptx
"""

import argparse
import os
import sys
import yaml
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree


# ═══════════════════════════════════════════════════════
# Layout Constants (LAYOUT_WIDE: 13.333" × 7.5")
# ═══════════════════════════════════════════════════════
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
HEADER_H = Inches(0.9)
MARGIN = Inches(0.5)
FONT = "Microsoft YaHei"


# ═══════════════════════════════════════════════════════
# Theme
# ═══════════════════════════════════════════════════════

def _parse_color(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class Theme:
    """Color palette driven by YAML config."""

    def __init__(self, cfg: dict | None = None):
        cfg = cfg or {}
        self.primary = _parse_color(cfg.get("primary", "#0A1F3D"))
        self.secondary = _parse_color(cfg.get("secondary", "#1A3A6B"))
        self.accent = _parse_color(cfg.get("accent", "#2196F3"))
        self.accent_light = _parse_color(cfg.get("accent_light", "#E3F2FD"))
        self.bg = _parse_color(cfg.get("bg", "#FFFFFF"))
        self.bg_light = _parse_color(cfg.get("bg_light", "#F5F7FA"))
        self.bg_card = _parse_color(cfg.get("bg_card", "#EEF2F7"))
        self.text_light = _parse_color(cfg.get("text_light", "#FFFFFF"))
        self.text_dark = _parse_color(cfg.get("text_dark", "#333333"))
        self.text_gray = _parse_color(cfg.get("text_gray", "#888888"))
        self.success = _parse_color(cfg.get("success", "#4CAF50"))
        self.warning = _parse_color(cfg.get("warning", "#FF9800"))
        self.meituan = _parse_color(cfg.get("meituan", "#FFC300"))
        self.douyin = _parse_color(cfg.get("douyin", "#252525"))
        self.font = cfg.get("font", FONT)


# ═══════════════════════════════════════════════════════
# Drawing Helpers
# ═══════════════════════════════════════════════════════

def _add_rect(slide, x, y, w, h, fill: RGBColor, line: RGBColor | None = None,
              rounded: bool = False, shadow: bool = False) -> object:
    """Add a (rounded) rectangle shape."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    if not shadow:
        shape.shadow.inherit = False
    return shape


def _add_text(slide, x, y, w, h, text: str,
              size: int = 14, color: RGBColor | None = None,
              bold: bool = False, align=PP_ALIGN.LEFT,
              font_name: str | None = None,
              anchor=MSO_ANCHOR.TOP) -> object:
    """Add a text box."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    # zero margins
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.alignment = align
    if color:
        p.font.color.rgb = color
    if font_name:
        p.font.name = font_name
    return tb


def _add_oval(slide, x, y, w, h, fill: RGBColor) -> object:
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _add_header(slide, title: str, theme: Theme, num: str = ""):
    prefix = f"{num}  " if num else ""
    full_title = prefix + title
    # Use title placeholder if available (from P25-style layout 2)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = full_title
            return
    # Fallback: add text box
    tb = _add_text(slide, MARGIN, Inches(0.25), Inches(12.3), Inches(0.55),
                   full_title, size=24, bold=True, font_name=theme.font,
                   anchor=MSO_ANCHOR.MIDDLE)
    run = tb.text_frame.paragraphs[0].runs[0]
    run.font.color.rgb = _parse_color("#4696EC")


# ═══════════════════════════════════════════════════════
# Slide Generators
# ═══════════════════════════════════════════════════════

def gen_agenda(slide, data: dict, theme: Theme):
    """4-section agenda: deep-blue full background with numbered list."""
    sections = data.get("sections", [])

    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme.primary)

    _add_text(slide, MARGIN, Inches(0.6), SLIDE_W - 2 * MARGIN, Inches(0.7),
              "目  录", size=30, color=theme.text_light, bold=True,
              align=PP_ALIGN.CENTER, font_name=theme.font)
    _add_rect(slide, Inches(5.5), Inches(1.35), Inches(2.3), Inches(0.04), theme.accent)

    y = Inches(2.0)
    for i, sec in enumerate(sections):
        _add_oval(slide, Inches(4.2), y, Inches(0.65), Inches(0.65), theme.accent)
        _add_text(slide, Inches(4.2), y, Inches(0.65), Inches(0.65),
                  str(i + 1), size=24, color=theme.text_light, bold=True,
                  align=PP_ALIGN.CENTER, font_name=theme.font, anchor=MSO_ANCHOR.MIDDLE)
        _add_text(slide, Inches(5.15), y, Inches(6), Inches(0.65),
                  sec, size=22, color=theme.text_light,
                  font_name=theme.font, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(1.05)


def gen_requirement_understanding(slide, data: dict, theme: Theme):
    """需求理解: 左侧分类卡片 + 右侧核心洞察"""
    _add_header(slide, data.get("title", "需求理解"), theme, "01")

    categories = data.get("categories", [])
    insights = data.get("key_insights", [])
    total = data.get("total", 0)

    # ── Left: category cards ──
    card_w = Inches(6.8)
    card_h = Inches(0.82)
    gap = Inches(0.14)
    y = Inches(1.25)

    for cat in categories:
        _add_rect(slide, MARGIN, y, card_w, card_h, theme.bg_light, rounded=True)
        _add_rect(slide, MARGIN, y, Inches(0.1), card_h, theme.accent)
        _add_text(slide, MARGIN + Inches(0.35), y + Inches(0.08), Inches(3.5), Inches(0.35),
                  cat.get("name", ""), size=15, color=theme.text_dark, bold=True,
                  font_name=theme.font)
        _add_text(slide, MARGIN + Inches(0.35), y + Inches(0.43), Inches(4.8), Inches(0.3),
                  " · ".join(cat.get("key_items", [])), size=11, color=theme.text_gray,
                  font_name=theme.font)
        # count badge
        badge_x = MARGIN + card_w - Inches(1.1)
        _add_rect(slide, badge_x, y + Inches(0.18), Inches(0.85), Inches(0.46),
                  theme.accent, rounded=True)
        _add_text(slide, badge_x, y + Inches(0.18), Inches(0.85), Inches(0.46),
                  f"{cat.get('count', 0)}项", size=14, color=theme.text_light, bold=True,
                  align=PP_ALIGN.CENTER, font_name=theme.font, anchor=MSO_ANCHOR.MIDDLE)
        y += card_h + gap

    # total
    _add_text(slide, MARGIN, Inches(6.6), card_w, Inches(0.4),
              f"合计 {total} 项功能需求", size=15, color=theme.secondary, bold=True,
              font_name=theme.font)

    # ── Right: key insights panel ──
    rx = Inches(7.6)
    rw = SLIDE_W - rx - MARGIN
    rh = Inches(5.5)
    _add_rect(slide, rx, Inches(1.25), rw, rh, theme.primary, rounded=True)

    _add_text(slide, rx + Inches(0.3), Inches(1.45), rw - Inches(0.6), Inches(0.5),
              "◆ 核心洞察", size=18, color=theme.text_light, bold=True, font_name=theme.font)

    iy = Inches(2.15)
    for insight in insights:
        _add_oval(slide, rx + Inches(0.3), iy + Inches(0.06), Inches(0.14), Inches(0.14),
                  theme.accent)
        _add_text(slide, rx + Inches(0.6), iy, rw - Inches(0.9), Inches(0.85),
                  insight, size=13, color=theme.text_light, font_name=theme.font)
        iy += Inches(0.95)


def gen_coverage_analysis(slide, data: dict, theme: Theme):
    """功能覆盖度: 三大指标卡居中"""
    _add_header(slide, data.get("title", "功能覆盖度分析"), theme, "02")

    stats = data.get("stats", {})
    sat = stats.get("satisfied", {})
    cust = stats.get("custom_dev", {})
    exc = stats.get("exceed", {})

    metrics = [
        {
            "big": f"{sat.get('pct', 85)}%",
            "label": sat.get("label", "直接满足"),
            "sub": f"{sat.get('value', 131)}项",
            "desc": "系统现有功能完全覆盖客户需求",
            "color": theme.success,
        },
        {
            "big": f"{cust.get('pct', 14)}%",
            "label": cust.get("label", "定制开发"),
            "sub": f"{cust.get('value', 22)}项",
            "desc": "少量定制开发即可满足",
            "color": theme.warning,
        },
        {
            "big": f"{exc.get('value', 141)}项",
            "label": exc.get("label", "超出期望"),
            "sub": "客户未要求",
            "desc": "蓝联功能超出客户需求范围",
            "color": theme.accent,
        },
    ]

    card_w = Inches(3.5)
    card_h = Inches(3.8)
    gap = Inches(0.5)
    total_w = card_w * 3 + gap * 2
    sx = (SLIDE_W - total_w) / 2

    for i, m in enumerate(metrics):
        x = sx + i * (card_w + gap)
        y = Inches(1.5)
        _add_rect(slide, x, y, card_w, card_h, theme.bg_light, rounded=True)
        _add_rect(slide, x, y, card_w, Inches(0.12), m["color"])
        # big number
        _add_text(slide, x, y + Inches(0.55), card_w, Inches(1.2),
                  m["big"], size=46, color=m["color"], bold=True,
                  align=PP_ALIGN.CENTER, font_name=theme.font)
        # label
        _add_text(slide, x, y + Inches(1.85), card_w, Inches(0.5),
                  m["label"], size=20, color=theme.text_dark, bold=True,
                  align=PP_ALIGN.CENTER, font_name=theme.font)
        # sub
        _add_text(slide, x, y + Inches(2.4), card_w, Inches(0.35),
                  m["sub"], size=14, color=theme.text_gray,
                  align=PP_ALIGN.CENTER, font_name=theme.font)
        # desc
        _add_text(slide, x + Inches(0.3), y + Inches(2.95), card_w - Inches(0.6), Inches(0.6),
                  m["desc"], size=12, color=theme.text_gray,
                  align=PP_ALIGN.CENTER, font_name=theme.font)

    # bottom statement
    _add_text(slide, MARGIN, Inches(6.0), SLIDE_W - 2 * MARGIN, Inches(0.5),
              "✓ 总体覆盖率 99%，系统成熟度高，可快速满足正祥商业会员系统建设需求",
              size=15, color=theme.secondary, bold=True,
              align=PP_ALIGN.CENTER, font_name=theme.font)


def gen_core_scenario(slide, data: dict, theme: Theme):
    """核心场景: 双行闭环流 (核心闭环版式)"""
    _add_header(slide, data.get("title", "核心场景"), theme, "03")

    loops = data.get("loops", [])
    y_start = Inches(1.5)
    row_h = Inches(2.4)

    for i, loop in enumerate(loops):
        y = y_start + i * row_h
        color_name = loop.get("color", "#2196F3")
        color = _parse_color(color_name)
        bg_color = _parse_color(loop.get("bg", "#FFFDE7" if i == 0 else "#F5F5F5"))

        # row background
        _add_rect(slide, MARGIN, y, SLIDE_W - 2 * MARGIN, Inches(2.0), bg_color, rounded=True)

        # loop name badge
        _add_rect(slide, MARGIN + Inches(0.2), y + Inches(0.25), Inches(1.8), Inches(0.55),
                  color, rounded=True)
        _add_text(slide, MARGIN + Inches(0.2), y + Inches(0.25), Inches(1.8), Inches(0.55),
                  loop.get("name", ""), size=15, color=theme.text_light, bold=True,
                  align=PP_ALIGN.CENTER, font_name=theme.font, anchor=MSO_ANCHOR.MIDDLE)

        # steps
        steps = loop.get("steps", [])
        n = len(steps)
        if n == 0:
            continue

        step_w = Inches(1.65)
        arrow_w = Inches(0.45)
        flow_w = step_w * n + arrow_w * (n - 1)
        flow_sx = MARGIN + Inches(2.3)
        flow_sy = y + Inches(0.6)

        for j, step in enumerate(steps):
            sx = flow_sx + j * (step_w + arrow_w)
            _add_rect(slide, sx, flow_sy, step_w, Inches(0.8), color, rounded=True)
            _add_text(slide, sx, flow_sy, step_w, Inches(0.8),
                      step, size=13, color=theme.text_light, bold=True,
                      align=PP_ALIGN.CENTER, font_name=theme.font,
                      anchor=MSO_ANCHOR.MIDDLE)
            if j < n - 1:
                _add_text(slide, sx + step_w, flow_sy, arrow_w, Inches(0.8),
                          "→", size=22, color=theme.text_gray, bold=True,
                          align=PP_ALIGN.CENTER, font_name=theme.font,
                          anchor=MSO_ANCHOR.MIDDLE)

        # loop-back indicator
        _add_text(slide, flow_sx, y + Inches(1.55), flow_w, Inches(0.35),
                  "↻ 闭环运营", size=11, color=theme.text_gray,
                  align=PP_ALIGN.CENTER, font_name=theme.font)


def gen_pricing_comparison(slide, data: dict, theme: Theme):
    """报价对比: 双方案大卡片"""
    _add_header(slide, data.get("title", "报价方案对比"), theme, "04")

    options = data.get("options", [])
    card_w = Inches(5.5)
    card_h = Inches(5.0)
    gap = Inches(0.8)
    total_w = card_w * len(options) + gap * (len(options) - 1)
    sx = (SLIDE_W - total_w) / 2
    y = Inches(1.3)

    for i, opt in enumerate(options):
        x = sx + i * (card_w + gap)
        recommended = opt.get("recommended", False)
        border = theme.accent if recommended else theme.bg_card

        # card
        _add_rect(slide, x, y, card_w, card_h, theme.bg_light, line=border, rounded=True)

        # header band
        hdr_color = theme.primary if recommended else theme.secondary
        _add_rect(slide, x, y, card_w, Inches(1.0), hdr_color, rounded=True)
        _add_text(slide, x + Inches(0.3), y + Inches(0.1), card_w - Inches(0.6), Inches(0.8),
                  f"{opt.get('icon', '')}  {opt.get('name', '')}", size=20,
                  color=theme.text_light, bold=True, font_name=theme.font,
                  anchor=MSO_ANCHOR.MIDDLE)

        # recommended badge
        if recommended:
            _add_rect(slide, x + card_w - Inches(1.2), y + Inches(0.15),
                      Inches(1.0), Inches(0.38), theme.accent, rounded=True)
            _add_text(slide, x + card_w - Inches(1.2), y + Inches(0.15),
                      Inches(1.0), Inches(0.38), "★ 推荐", size=12,
                      color=theme.text_light, bold=True,
                      align=PP_ALIGN.CENTER, font_name=theme.font,
                      anchor=MSO_ANCHOR.MIDDLE)

        # price
        _add_text(slide, x + Inches(0.4), y + Inches(1.2), card_w - Inches(0.8), Inches(0.35),
                  "首年费用", size=12, color=theme.text_gray, font_name=theme.font)
        _add_text(slide, x + Inches(0.4), y + Inches(1.55), card_w - Inches(0.8), Inches(0.6),
                  opt.get("first_year", ""), size=28, color=theme.primary, bold=True,
                  font_name=theme.font)
        _add_text(slide, x + Inches(0.4), y + Inches(2.25), card_w - Inches(0.8), Inches(0.35),
                  f"次年续费：{opt.get('renewal', '')}", size=13, color=theme.text_gray,
                  font_name=theme.font)

        # features
        fy = y + Inches(2.9)
        _add_text(slide, x + Inches(0.4), fy, card_w - Inches(0.8), Inches(0.35),
                  "方案特点", size=14, color=theme.text_dark, bold=True, font_name=theme.font)
        fy += Inches(0.42)
        for feat in opt.get("features", []):
            _add_text(slide, x + Inches(0.55), fy, card_w - Inches(1.0), Inches(0.32),
                      f"✓  {feat}", size=12, color=theme.text_gray, font_name=theme.font)
            fy += Inches(0.36)


def gen_implementation_plan(slide, data: dict, theme: Theme):
    """实施计划: 阶段时间线"""
    _add_header(slide, data.get("title", "实施计划"), theme, "05")

    phases = data.get("phases", [])
    if not phases:
        return

    pw = Inches(5.5)
    ph = Inches(4.5)
    gap = Inches(0.8)
    total_w = pw * len(phases) + gap * (len(phases) - 1)
    sx = (SLIDE_W - total_w) / 2
    y = Inches(1.4)

    for i, phase in enumerate(phases):
        x = sx + i * (pw + gap)

        _add_rect(slide, x, y, pw, ph, theme.bg_light, rounded=True)
        _add_rect(slide, x, y, pw, Inches(0.1), theme.accent)

        # number circle
        _add_oval(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.65), Inches(0.65),
                  theme.accent)
        _add_text(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.65), Inches(0.65),
                  str(i + 1), size=22, color=theme.text_light, bold=True,
                  align=PP_ALIGN.CENTER, font_name=theme.font, anchor=MSO_ANCHOR.MIDDLE)

        # name + duration
        _add_text(slide, x + Inches(1.15), y + Inches(0.28), Inches(4), Inches(0.35),
                  phase.get("name", ""), size=16, color=theme.primary, bold=True,
                  font_name=theme.font)
        _add_text(slide, x + Inches(1.15), y + Inches(0.62), Inches(4), Inches(0.3),
                  phase.get("duration", ""), size=13, color=theme.accent,
                  font_name=theme.font)

        # description
        _add_text(slide, x + Inches(0.3), y + Inches(1.2), pw - Inches(0.6), Inches(0.45),
                  phase.get("desc", ""), size=15, color=theme.text_dark, bold=True,
                  font_name=theme.font)

        # items
        iy = y + Inches(1.85)
        for item in phase.get("items", []):
            _add_oval(slide, x + Inches(0.35), iy + Inches(0.1), Inches(0.1), Inches(0.1),
                      theme.accent)
            _add_text(slide, x + Inches(0.6), iy, pw - Inches(0.9), Inches(0.35),
                      item, size=13, color=theme.text_gray, font_name=theme.font)
            iy += Inches(0.4)

        # arrow between phases
        if i < len(phases) - 1:
            ax = x + pw
            _add_text(slide, ax, y + Inches(1.5), gap, Inches(1),
                      "→", size=36, color=theme.accent, bold=True,
                      align=PP_ALIGN.CENTER, font_name=theme.font,
                      anchor=MSO_ANCHOR.MIDDLE)

    # bottom note
    note = data.get("note", "")
    if note:
        _add_text(slide, MARGIN, Inches(6.4), SLIDE_W - 2 * MARGIN, Inches(0.4),
                  note, size=13, color=theme.text_gray,
                  align=PP_ALIGN.CENTER, font_name=theme.font)


def gen_text_bullets(slide, data: dict, theme: Theme):
    """Generic text+bullets slide for company intro mode."""
    _add_header(slide, data.get("title", ""), theme)
    body = data.get("body", "")
    items = data.get("items", [])

    y = Inches(1.3)
    if body:
        _add_text(slide, MARGIN, y, SLIDE_W - 2 * MARGIN, Inches(1.2),
                  body, size=15, color=theme.text_dark, font_name=theme.font)
        y += Inches(1.3)

    for item in items:
        _add_oval(slide, MARGIN + Inches(0.1), y + Inches(0.12), Inches(0.1), Inches(0.1),
                  theme.accent)
        text = item if isinstance(item, str) else item.get("text", "")
        _add_text(slide, MARGIN + Inches(0.4), y, SLIDE_W - 2 * MARGIN - Inches(0.4), Inches(0.45),
                  text, size=14, color=theme.text_dark, font_name=theme.font)
        y += Inches(0.55)


def gen_feature_cards(slide, data: dict, theme: Theme):
    """Feature card grid for company intro mode."""
    _add_header(slide, data.get("title", ""), theme)
    body = data.get("body", "")
    items = data.get("items", [])
    cols = data.get("columns", 3)

    y = Inches(1.3)
    if body:
        _add_text(slide, MARGIN, y, SLIDE_W - 2 * MARGIN, Inches(0.5),
                  body, size=14, color=theme.text_gray, font_name=theme.font)
        y += Inches(0.65)

    if not items:
        return

    gap = Inches(0.3)
    card_w = (SLIDE_W - 2 * MARGIN - gap * (cols - 1)) / cols
    card_h = Inches(2.2)
    rows = (len(items) + cols - 1) // cols

    for i, item in enumerate(items):
        col = i % cols
        row = i // cols
        x = MARGIN + col * (card_w + gap)
        cy = y + row * (card_h + gap)

        _add_rect(slide, x, cy, card_w, card_h, theme.bg_light, rounded=True)
        _add_rect(slide, x, cy, card_w, Inches(0.08), theme.accent)
        _add_text(slide, x + Inches(0.25), cy + Inches(0.2), card_w - Inches(0.5), Inches(0.4),
                  item.get("title", ""), size=15, color=theme.primary, bold=True,
                  font_name=theme.font)
        _add_text(slide, x + Inches(0.25), cy + Inches(0.7), card_w - Inches(0.5), Inches(1.3),
                  item.get("body", ""), size=12, color=theme.text_gray, font_name=theme.font)


# ═══════════════════════════════════════════════════════
# Slide Registry
# ═══════════════════════════════════════════════════════

SLIDE_GENERATORS = {
    "agenda": gen_agenda,
    "requirement-understanding": gen_requirement_understanding,
    "coverage-analysis": gen_coverage_analysis,
    "core-scenario": gen_core_scenario,
    "pricing-comparison": gen_pricing_comparison,
    "implementation-plan": gen_implementation_plan,
    "text-bullets": gen_text_bullets,
    "feature-cards": gen_feature_cards,
}


# ═══════════════════════════════════════════════════════
# Cover Modifier
# ═══════════════════════════════════════════════════════

def modify_why_choose_slide(slide):
    A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if "八大优势" not in shape.text_frame.text:
            continue
        if shape.top and shape.top > Emu(914400):
            continue
        runs = shape._element.findall(f".//{{{A_NS}}}r")
        if not runs:
            continue
        first_t = runs[0].find(f"{{{A_NS}}}t")
        if first_t is not None:
            first_t.text = "为什么选择蓝联 — 八大核心优势 + 运营营销双轮驱动"
        for r in runs[1:]:
            r.getparent().remove(r)
        return True
    return False


def modify_cover(slide, cover_data: dict, theme: Theme):
    """Modify existing cover slide: add client name text box."""
    if not cover_data:
        return

    title = cover_data.get("title", "")
    if title:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text.strip()
                if "商圈会员数字营销解决方案" in text:
                    shape.text_frame.paragraphs[0].text = title
                    break

    # Add client name below existing title area
    subtitle = cover_data.get("subtitle", "")
    if subtitle:
        _add_text(slide, Inches(2.5), Inches(5.2), Inches(8.3), Inches(0.6),
                  subtitle, size=18, color=theme.accent,
                  align=PP_ALIGN.CENTER, font_name=theme.font)

    date_str = cover_data.get("date", "")
    if date_str:
        _add_text(slide, Inches(2.5), Inches(5.9), Inches(8.3), Inches(0.4),
                  f"方案日期：{date_str}", size=13, color=theme.text_gray,
                  align=PP_ALIGN.CENTER, font_name=theme.font)


GRAD_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _set_run_gradient(run, c1: str, c2: str):
    rPr = run._r.get_or_add_rPr()
    for tag in ("gradFill", "solidFill", "noFill"):
        for elem in rPr.findall(f"{{{GRAD_NS}}}{tag}"):
            rPr.remove(elem)
    grad = etree.fromstring(
        f'<a:gradFill xmlns:a="{GRAD_NS}">'
        f'<a:gsLst>'
        f'<a:gs pos="44000"><a:srgbClr val="{c1}"/></a:gs>'
        f'<a:gs pos="100000"><a:srgbClr val="{c2}"/></a:gs>'
        f'</a:gsLst>'
        f'<a:lin ang="5400000" scaled="0"/>'
        f'</a:gradFill>'
    )
    rPr.insert(0, grad)


def _set_section_text(shape, num: int, text: str, active: bool = False):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"{num}. {text}"
    p.alignment = PP_ALIGN.LEFT
    if not p.runs:
        p.add_run()
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.name = "微软雅黑"
    run.font.italic = active
    if active:
        _set_run_gradient(run, "2C5DE6", "6F92F3")
    else:
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


TOC_SLIDE_INFO = {
    1: 1,
    36: 2,
    9: 3,
    55: 4,
}


def modify_toc_template_slide(slide, sections: list[str], active_section: int = 0):
    if not sections:
        return

    text_shapes = [s for s in slide.shapes if getattr(s, "has_text_frame", False)]
    text_shapes.sort(key=lambda s: s.top)
    if len(text_shapes) < 3:
        return

    existing = len(text_shapes)
    n = min(len(sections), 4)

    if existing >= n:
        for i in range(n):
            _set_section_text(text_shapes[i], i + 1, sections[i],
                              active=(i + 1 == active_section))
    else:
        for i in range(existing):
            _set_section_text(text_shapes[i], i + 1, sections[i],
                              active=(i + 1 == active_section))
        x = text_shapes[0].left
        y_step = text_shapes[1].top - text_shapes[0].top
        last_top = text_shapes[existing - 1].top
        w = max(s.width for s in text_shapes[:existing])
        h = text_shapes[0].height
        for i in range(existing, n):
            y = last_top + y_step * (i - existing + 1)
            tb = slide.shapes.add_textbox(x, y, w, h)
            _set_section_text(tb, i + 1, sections[i],
                              active=(i + 1 == active_section))


def fix_sidebar_order(slide, desired_order: list[str]):
    cat_set = set(desired_order)
    cat_shapes: dict[str, object] = {}
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        txt = shape.text_frame.text.strip()
        if txt in cat_set and shape.left and shape.left < Emu(914400):
            cat_shapes[txt] = shape
    if len(cat_shapes) < len(desired_order):
        return False
    y_slots = sorted(s.top for s in cat_shapes.values())
    for i, cat in enumerate(desired_order):
        if cat in cat_shapes:
            cat_shapes[cat].top = y_slots[i]
    return True


# ═══════════════════════════════════════════════════════
# Slide Manipulation
# ═══════════════════════════════════════════════════════

def move_slides_after_position(prs, count: int, target_pos: int = 1):
    """Move the last `count` slides to start at `target_pos`."""
    xml = prs.slides._sldIdLst
    elements = list(xml)
    total = len(elements)

    for i in range(count):
        last = elements[total - count + i]
        xml.remove(last)
        xml.insert(target_pos + i, last)


# Reference PPT (蓝联科技CRM商圈会员数智营销方案_202604.pptx) product detail mapping
# 0-indexed slide → category, covering slides 15-34 (indices 14-33)
PRODUCT_SLIDE_CATEGORIES: dict[int, str] = {
    14: "会员",   # S15 多元化的会员服务
    15: "积分",   # S16 配套完善的积分体系
    16: "积分",   # S17 支付即积分场景
    17: "积分",   # S18 多元积分生态
    18: "积分",   # S19 强大的积分管理体系
    19: "营销",   # S20 场景化营销闭环
    20: "营销",   # S21 丰富的营销工具
    21: "营销",   # S22 新媒体多域内容营销
    22: "运营",   # S23 种草社区
    23: "商城",   # S24 电商平台4+1商城
    24: "商城",   # S25 门店自提
    25: "服务",   # S26 商户助手小程序
    26: "服务",   # S27 寻店寻优惠服务
    27: "服务",   # S28 停车服务
    28: "财务",   # S29 财务收银对账分账
    29: "财务",   # S30
    30: "财务",   # S31
    31: "财务",   # S32 支付宝碰一下会员
    32: "智能",   # S33 AI客服
    33: "智能",   # S34 商业智能分析平台
}


def reorder_product_slides(prs, category_order: list[str]):
    """Reorder product detail slides by category.

    Operates on the 0-indexed PRODUCT_SLIDE_CATEGORIES mapping.
    Only affects slides within that range; all others stay in place.
    """
    indices = sorted(PRODUCT_SLIDE_CATEGORIES.keys())
    groups: dict[str, list[int]] = {}
    for idx in indices:
        cat = PRODUCT_SLIDE_CATEGORIES[idx]
        groups.setdefault(cat, []).append(idx)

    new_order: list[int] = []
    for cat in category_order:
        new_order.extend(groups.get(cat, []))
    for idx in indices:
        if idx not in new_order:
            new_order.append(idx)

    if new_order == indices:
        return  # no change needed

    xml = prs.slides._sldIdLst
    elements = list(xml)
    start = indices[0]

    product_elems = [elements[idx] for idx in indices]
    for elem in product_elems:
        xml.remove(elem)
    for i, old_idx in enumerate(new_order):
        xml.insert(start + i, elements[old_idx])


# ═══════════════════════════════════════════════════════
# Full-Deck Section Reorder
# ═══════════════════════════════════════════════════════

def _build_product_indices(category_order: list[str],
                          within_order: dict | None = None) -> list[int]:
    groups: dict[str, list[int]] = {}
    for idx in sorted(PRODUCT_SLIDE_CATEGORIES.keys()):
        cat = PRODUCT_SLIDE_CATEGORIES[idx]
        groups.setdefault(cat, []).append(idx)
    result: list[int] = []
    within_order = within_order or {}
    for cat in category_order:
        indices = groups.get(cat, [])
        if cat in within_order:
            ov = within_order[cat]
            indices = [i for i in ov if i in indices] + [i for i in indices if i not in ov]
        result.extend(indices)
    for idx in sorted(PRODUCT_SLIDE_CATEGORIES.keys()):
        if idx not in result:
            result.append(idx)
    return result


def build_section_deck_order(
    num_base: int,
    num_custom: int,
    product_categories: list[str] | None = None,
    within_order: dict | None = None,
) -> list[int]:
    cs = num_base

    cat_order = product_categories or [
        "会员", "积分", "服务", "运营", "营销", "商城", "财务", "智能"
    ]
    prod = _build_product_indices(cat_order, within_order)

    section1 = list(range(37, 55))
    # Custom slides: first 3 go to section 2, remaining go to section 4
    custom_in_s2 = min(3, num_custom)
    section2 = [cs + i for i in range(custom_in_s2)]
    custom_in_s4 = [cs + i for i in range(custom_in_s2, num_custom)]
    section3_pre = [5, 9, 10, 11, 12, 13]
    section3_post = [34]
    section3 = section3_pre + prod + section3_post
    section4 = [55] + custom_in_s4 + [56, 57, 58, 59, 60, 35, 61]

    return [0, 1] + section1 + [36] + section2 + section3 + section4


def build_intro_deck_order(
    num_base: int,
    section_counts: list[int],
) -> list[int]:
    cs = num_base
    section1 = list(range(37, 55))
    toc_indices = [1, 36, 9, 55]

    result = [0, toc_indices[0]] + section1
    offset = cs
    for i, count in enumerate(section_counts):
        sec_slides = [offset + j for j in range(count)]
        result += [toc_indices[i + 1]] + sec_slides
        offset += count
    result += [61]
    return result


def build_tender_deck_order(
    num_base: int,
    section_counts: list[int],
) -> list[int]:
    cs = num_base
    toc_indices = [1, 36, 9, 55]

    result = [0]
    offset = cs
    for i, count in enumerate(section_counts):
        toc_idx = toc_indices[i % len(toc_indices)]
        sec_slides = [offset + j for j in range(count)]
        result += [toc_idx] + sec_slides
        offset += count
    result += [61]
    return result


def apply_full_reorder(prs, final_order: list[int]):
    """Reorder ALL slides to match final_order. Leftovers go to end."""
    xml = prs.slides._sldIdLst
    elements = list(xml)
    total = len(elements)

    ordered = [elements[i] for i in final_order if i < total]
    in_order = set(i for i in final_order if i < total)
    leftovers = [elements[i] for i in range(total) if i not in in_order]

    for elem in elements:
        xml.remove(elem)
    for elem in ordered + leftovers:
        xml.append(elem)


def _drop_unused_slides(prs, keep: int):
    xml = prs.slides._sldIdLst
    elements = list(xml)
    for elem in elements[keep:]:
        rId = elem.get(qn("r:id"))
        xml.remove(elem)
        if rId:
            prs.part.drop_rel(rId)


def _set_slide_layout(slide, new_layout):
    """Change slide layout via relationship update. new_layout from prs.slide_layouts[N]."""
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    slide_part = slide.part
    layout_part = new_layout.part
    for rel in list(slide_part.rels.values()):
        if rel.reltype == RT.SLIDE_LAYOUT:
            rel._target = layout_part
            return True
    return False


# ═══════════════════════════════════════════════════════
# Main Compiler
# ═══════════════════════════════════════════════════════

def compile_proposal(content_path: str, base_path: str | None = None,
                     output_path: str | None = None) -> str:
    with open(content_path, "r", encoding="utf-8") as f:
        content = yaml.safe_load(f)

    base_path = base_path or content.get("base_ppt", "")
    output_path = output_path or content.get("output", "output.pptx")
    theme = Theme(content.get("theme", {}))

    if not base_path or not os.path.exists(base_path):
        print(f"❌ Base PPT not found: {base_path}", file=sys.stderr)
        sys.exit(1)

    # ── Step 1: Open base PPT ──
    prs = Presentation(base_path)
    num_base = len(prs.slides)
    mode = content.get("mode", "proposal")
    print(f"📖 Base PPT: {base_path} ({num_base} slides, mode={mode})")

    # ── Step 2: Modify cover ──
    cover_data = content.get("cover", {})
    if cover_data:
        modify_cover(prs.slides[0], cover_data, theme)
        print(f"   ✎ Cover: {cover_data.get('subtitle', '')}")

    if mode == "proposal":
        if modify_why_choose_slide(prs.slides[35]):
            print("   ✎ P35 title → 为什么选择蓝联")

        # ── Step 2.5: Standardize slide layouts ──
        # Convention: 首尾页 → 标题幻灯片(0), 目录页 → 节标题(1), 其他 → 标题和内容(2)
        content_layout = prs.slide_layouts[2]  # '标题和内容' — unified title+content
        for idx in range(56, 61):  # P53-P57: section 4 implementation info slides
            _set_slide_layout(prs.slides[idx], content_layout)
        print("   ↻ Layout: content slides → 标题和内容")

    # ── Step 3: Generate ALL custom slides at end ──
    title_layout = prs.slide_layouts[2]
    num_custom = 0
    section_counts: list[int] = []

    if mode in ("intro", "tender"):
        sections_cfg = content.get("sections", [])
        section_counts = []
        for sec in sections_cfg:
            sec_slides = sec.get("slides", [])
            count = 0
            for cfg in sec_slides:
                gen_func = SLIDE_GENERATORS.get(cfg.get("type", ""))
                if not gen_func:
                    print(f"   ⚠ Unknown type: {cfg.get('type', '')}")
                    continue
                slide = prs.slides.add_slide(title_layout)
                gen_func(slide, cfg, theme)
                count += 1
                num_custom += 1
                print(f"   + {cfg.get('type', '')}")
            section_counts.append(count)
    else:
        slides_config = content.get("slides", [])
        for cfg in slides_config:
            slide_type = cfg.get("type", "")
            gen_func = SLIDE_GENERATORS.get(slide_type)
            if not gen_func:
                print(f"   ⚠ Unknown type: {slide_type}")
                continue
            slide = prs.slides.add_slide(title_layout)
            gen_func(slide, cfg, theme)
            num_custom += 1
            print(f"   + {slide_type}")

    # ── Step 4: Reorder ──
    if mode == "intro":
        toc_titles = content.get("toc", {}).get("titles", [])
        toc_slide_indices = [1, 36, 9, 55]
        for i, toc_idx in enumerate(toc_slide_indices):
            if toc_idx < num_base and toc_titles:
                modify_toc_template_slide(prs.slides[toc_idx], toc_titles, i + 1)
        final_order = build_intro_deck_order(num_base, section_counts)
        apply_full_reorder(prs, final_order)
        _drop_unused_slides(prs, len(final_order))
        print(f"   ↻ Intro reorder: cover → company+cases → {len(section_counts)} custom sections → closing")
    elif mode == "tender":
        toc_titles = content.get("toc", {}).get("titles", [])
        toc_slide_indices = [1, 36, 9, 55]
        for i, toc_idx in enumerate(toc_slide_indices):
            if toc_idx < num_base and toc_titles:
                modify_toc_template_slide(prs.slides[toc_idx], toc_titles, i + 1)
        final_order = build_tender_deck_order(num_base, section_counts)
        apply_full_reorder(prs, final_order)
        _drop_unused_slides(prs, len(final_order))
        print(f"   ↻ Tender reorder: cover → {len(section_counts)} tender sections → closing")
    elif sections_cfg := content.get("sections"):
        if sections_cfg.get("enabled", True):
            titles = sections_cfg.get("titles", [])
            for slide_idx, active_sec in TOC_SLIDE_INFO.items():
                if slide_idx < num_base:
                    modify_toc_template_slide(prs.slides[slide_idx], titles, active_sec)
            cat_order = content.get("reorder", {}).get("order")
            within_order = content.get("reorder", {}).get("within")
            if cat_order:
                fixed = 0
                for idx in sorted(PRODUCT_SLIDE_CATEGORIES.keys()):
                    if fix_sidebar_order(prs.slides[idx], cat_order):
                        fixed += 1
                print(f"   ↻ Sidebar fixed on {fixed} product slides")
            final_order = build_section_deck_order(num_base, num_custom, cat_order, within_order)
            apply_full_reorder(prs, final_order)
            _drop_unused_slides(prs, len(final_order))
            print(f"   ↻ 4-section reorder: {' → '.join(titles)}")
    else:
        reorder_cfg = content.get("reorder")
        if reorder_cfg and reorder_cfg.get("order"):
            reorder_product_slides(prs, reorder_cfg["order"])
            print(f"   ↻ Product reorder: {' → '.join(reorder_cfg['order'])}")
        if num_custom > 0:
            insert_after = content.get("insert_after", 1)
            move_slides_after_position(prs, num_custom, insert_after)
            print(f"   ↻ {num_custom} custom slides → position {insert_after + 1}")

    # ── Step 5: Save ──
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    size_kb = os.path.getsize(output_path) // 1024
    print(f"\n✅ Output: {output_path}")
    print(f"📊 Total: {len(prs.slides)} slides ({size_kb}KB)")

    return output_path


# ═══════════════════════════════════════════════════════
# CLI
# ═══════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(
        description="方案介绍PPT模板编译器 — 基于参考PPT + YAML内容包生成客户定制方案"
    )
    parser.add_argument("content", help="YAML content package path")
    parser.add_argument("--base", default=None, help="Base/reference PPTX (overrides YAML)")
    parser.add_argument("--output", default=None, help="Output PPTX path (overrides YAML)")
    args = parser.parse_args()

    compile_proposal(args.content, args.base, args.output)


if __name__ == "__main__":
    main()
