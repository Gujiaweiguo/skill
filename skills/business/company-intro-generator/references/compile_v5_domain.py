#!/usr/bin/env python3
"""
V5 领域驱动方案编译器 — YAML 驱动零代码生成 PPT

用法:
    python compile_v5_domain.py <config.yaml>

配置规范见: domain-slide-config-schema.yaml

特性:
    - 完全 YAML 驱动，零代码生成 48+ 页方案 PPT
    - 21 种 slide type 覆盖需求理解/蓝图/领域展开/实施/结尾全流程
    - 可覆盖调色板与几何参数
    - 复用 base PPT 公司介绍+案例段，自定义页无缝拼接
    - 支持 $PROPOSALS_DIR / $LANLNK_BASE 环境变量展开

依赖:
    pip install python-pptx pyyaml lxml
"""

import os
import sys
import math
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# =============================================================================
# Default Palette & Geometry
# =============================================================================

DEFAULT_PALETTE = {
    'ink':        '#18243B',
    'body':       '#3D4A5C',
    'muted':      '#8A94A6',
    'blue':       '#2C5DE6',
    'blue_dk':    '#1E47B8',
    'blue_lt':    '#D8E4FA',
    'blue_bg':    '#F5F8FE',
    'green':      '#0EA56F',
    'teal':       '#61C9BD',
    'amber':      '#F59E1B',
    'orange':     '#E06C00',
    'line':       '#E2E8F0',
    'white':      '#FFFFFF',
    'gray_light': '#F5F5F5',
    'green_light':'#ECF9F4',
    'gray_bg':    '#F8F8F8',
    'font':       'Microsoft YaHei',
    'title_font': 'Alibaba PuHuiTi Bold',
}

DEFAULT_GEOMETRY = {
    'page_w': 13.333,
    'tx': 0.6, 'tw': 12.1,
    'lx': 0.6, 'lw': 7.6,
    'rx': 8.6, 'rw': 4.1,
    'cy': 1.2,
}


def _hex_to_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _expand_path(path):
    """Expand environment variables in a path."""
    if path is None:
        return None
    return os.path.expandvars(path)


# =============================================================================
# Renderer — drawing context with palette + geometry
# =============================================================================

class Renderer:
    """Holds palette/geometry and provides all drawing primitives."""

    def __init__(self, palette_cfg=None, geometry_cfg=None):
        # Build palette: start with defaults (as RGBColor), override from config
        self.pal = {}
        for k, v in DEFAULT_PALETTE.items():
            if k in ('font', 'title_font'):
                self.pal[k] = v
            else:
                self.pal[k] = _hex_to_rgb(v)
        if palette_cfg:
            for k, v in palette_cfg.items():
                if k in ('font', 'title_font'):
                    self.pal[k] = v
                elif isinstance(v, str):
                    self.pal[k] = _hex_to_rgb(v) if v.startswith('#') else self.pal.get(k, _hex_to_rgb(DEFAULT_PALETTE.get(k, '#000000')))

        # Build geometry
        g = dict(DEFAULT_GEOMETRY)
        if geometry_cfg:
            g.update(geometry_cfg)
        self.page_w = Inches(g['page_w'])
        self.tx = Inches(g['tx']);  self.tw = Inches(g['tw'])
        self.lx = Inches(g['lx']);  self.lw = Inches(g['lw'])
        self.rx = Inches(g['rx']);  self.rw = Inches(g['rw'])
        self.cy = Inches(g['cy'])

    def color(self, name):
        """Resolve a color name/hex to RGBColor."""
        if name is None:
            return None
        if isinstance(name, RGBColor):
            return name
        if isinstance(name, str) and name.startswith('#'):
            return _hex_to_rgb(name)
        return self.pal.get(name, self.pal['blue'])

    # ── Basic primitives ──

    def txt(self, s, x, y, w, h, text, sz=16, color='body', bold=False,
            align=PP_ALIGN.LEFT, font=None, anchor=MSO_ANCHOR.TOP, spacing=4):
        tb = s.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, m, 0)
        c = self.color(color) if color else None
        f = font or self.pal['font']
        lines = text.split("\n")
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = ln; p.font.size = Pt(sz); p.font.bold = bold
            p.font.name = f; p.alignment = align; p.space_after = Pt(spacing)
            if c: p.font.color.rgb = c
        return tb

    def line(self, s, x, y, w, color='line', weight=0.75):
        ln = s.shapes.add_connector(1, x, y, x + w, y)
        ln.line.color.rgb = self.color(color); ln.line.width = Pt(weight)

    def dot(self, s, x, y, d=Inches(0.1), color='blue'):
        o = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
        o.fill.solid(); o.fill.fore_color.rgb = self.color(color)
        o.line.fill.background(); o.shadow.inherit = False
        return o

    def rect(self, s, x, y, w, h, fill, rounded=False):
        st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        r = s.shapes.add_shape(st, x, y, w, h)
        r.fill.solid(); r.fill.fore_color.rgb = self.color(fill)
        r.line.fill.background(); r.shadow.inherit = False
        return r

    def rrect(self, s, x, y, w, h, fill, line_color=None):
        r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        r.fill.solid(); r.fill.fore_color.rgb = self.color(fill)
        if line_color:
            r.line.color.rgb = self.color(line_color); r.line.width = Pt(0.75)
        else:
            r.line.fill.background()
        r.shadow.inherit = False
        return r

    # ── Page scaffolding ──

    def page_title(self, s, title):
        self.txt(s, self.tx, Inches(0.12), self.tw, Inches(0.45),
                 title, sz=22, color='ink', bold=True, font=self.pal['title_font'],
                 anchor=MSO_ANCHOR.MIDDLE)

    def page_source(self, s, text):
        if text:
            self.txt(s, self.tx, Inches(7.12), Inches(10), Inches(0.25),
                     text, sz=10, color='muted')

    def so_what(self, s, bullets, heading="关键结论"):
        rx, rw, cy = self.rx, self.rw, self.cy
        self.rect(s, rx - Inches(0.12), cy - Inches(0.08), rw + Inches(0.24), Inches(5.5), 'blue_bg')
        self.rect(s, rx - Inches(0.12), cy - Inches(0.08), Inches(0.05), Inches(5.5), 'blue')
        self.txt(s, rx, cy + Inches(0.05), rw, Inches(0.4), heading,
                 sz=17, color='blue_dk', bold=True, font=self.pal['title_font'])
        y = cy + Inches(0.65)
        for b in bullets:
            self.rrect(s, rx, y + Inches(0.04), Inches(0.16), Inches(0.16), 'blue_lt')
            self.txt(s, rx + Inches(0.28), y, rw - Inches(0.28), Inches(0.7), b,
                     sz=15, color='body', spacing=2)
            y += Inches(0.82)

    def big_number(self, s, x, y, w, num, label, color='blue'):
        self.txt(s, x, y, w, Inches(0.85), num, sz=48, color=color, bold=True,
                 font=self.pal['title_font'], align=PP_ALIGN.LEFT)
        self.txt(s, x, y + Inches(0.78), w, Inches(0.35), label, sz=15, color='muted')

    def bottom_bar(self, s, text, color='blue'):
        self.rect(s, self.tx, Inches(6.35), Inches(0.06), Inches(0.5), color)
        self.txt(s, self.tx + Inches(0.18), Inches(6.35), self.tw - Inches(0.18), Inches(0.5),
                 text, sz=15, color='ink', bold=True, anchor=MSO_ANCHOR.MIDDLE)


# =============================================================================
# Slide Builders — one function per type
# =============================================================================

def build_requirement_projects(s, cfg, r):
    r.page_title(s, cfg['title'])
    lx, lw, cy = r.lx, r.lw, r.cy
    pilot = cfg['pilot']
    r.txt(s, lx, cy, Inches(3.5), Inches(0.35), "首期试点", sz=13, color='blue', bold=True, font=r.pal['title_font'])
    r.txt(s, lx, cy + Inches(0.35), Inches(3.5), Inches(0.5), pilot['name'], sz=22, color='ink', bold=True, font=r.pal['title_font'])
    r.txt(s, lx, cy + Inches(0.9), Inches(3.5), Inches(0.35), f"{pilot['type']} · {pilot['desc']}", sz=15, color='body')
    r.line(s, lx, cy + Inches(1.35), Inches(3.2))
    r.txt(s, lx, cy + Inches(1.5), Inches(3.5), Inches(2), pilot['detail'], sz=14, color='muted', spacing=4)
    ox = lx + Inches(4.0)
    r.txt(s, ox, cy, Inches(3.5), Inches(0.35), cfg.get('phase2_label', '二期/三期'), sz=13, color='muted', bold=True, font=r.pal['title_font'])
    for i, p in enumerate(cfg['others']):
        y = cy + Inches(0.45) + i * Inches(1.15)
        r.txt(s, ox, y, Inches(3.5), Inches(0.3), p['name'], sz=17, color='ink', bold=True, font=r.pal['title_font'])
        r.txt(s, ox, y + Inches(0.32), Inches(3.5), Inches(0.25), f"{p['type']} · {p['desc']}", sz=13, color='body')
    r.so_what(s, cfg.get('so_what', []))
    r.page_source(s, cfg.get('source'))


def build_requirement_roles(s, cfg, r):
    r.page_title(s, cfg['title'])
    lx, cy = r.lx, r.cy
    roles = cfg['roles']
    for i, role in enumerate(roles):
        col = i % 2; row = i // 2
        x = lx + col * Inches(3.8)
        y = cy + Inches(0.1) + row * Inches(1.6)
        r.txt(s, x, y, Inches(3.5), Inches(0.35), role['name'], sz=18, color='blue', bold=True, font=r.pal['title_font'])
        r.line(s, x, y + Inches(0.38), Inches(0.5), 'blue', 1.5)
        r.txt(s, x, y + Inches(0.5), Inches(3.5), Inches(1.0), role['pain'], sz=15, color='body', spacing=3)
    r.so_what(s, cfg.get('so_what', []))
    r.page_source(s, cfg.get('source'))


def build_requirement_pain_detail(s, cfg, r):
    r.page_title(s, cfg['title'])
    lx, cy = r.lx, r.cy
    for i, cat in enumerate(cfg['categories']):
        x = lx + i * Inches(2.55)
        r.txt(s, x, cy, Inches(2.3), Inches(0.35), cat['role'], sz=17, color='blue', bold=True, font=r.pal['title_font'])
        r.line(s, x, cy + Inches(0.35), Inches(2.2))
        for j, item in enumerate(cat['items']):
            if not item: continue
            y = cy + Inches(0.5) + j * Inches(0.55)
            r.dot(s, x, y + Inches(0.08), Inches(0.07))
            r.txt(s, x + Inches(0.18), y, Inches(2.1), Inches(0.45), item, sz=14, color='body', spacing=1)
    r.so_what(s, cfg.get('so_what', []))
    r.page_source(s, cfg.get('source'))


def build_requirement_matrix(s, cfg, r):
    r.page_title(s, cfg['title'])
    lx, cy = r.lx, r.cy
    headers = cfg['headers']
    cw = Inches(0.82); lh = Inches(1.3)
    for j, h in enumerate(headers):
        r.txt(s, lx + lh + j * cw, cy, cw, Inches(0.3), h, sz=12, color='muted', bold=True, align=PP_ALIGN.CENTER)
    for i, row in enumerate(cfg['rows']):
        y = cy + Inches(0.4) + i * Inches(0.5)
        r.txt(s, lx, y, lh, Inches(0.35), row['label'], sz=14, color='ink', bold=True)
        for j, v in enumerate(row['values']):
            stars = "●" * v + "○" * (5 - v)
            c = 'blue' if v >= 4 else 'muted'
            r.txt(s, lx + lh + j * cw, y, cw, Inches(0.35), stars, sz=10, color=c, align=PP_ALIGN.CENTER)
    legend = cfg.get('legend', "● 需求强    ○ 需求弱")
    r.txt(s, lx, cy + Inches(3.2), Inches(7), Inches(0.3), legend, sz=11, color='muted')
    r.so_what(s, cfg.get('so_what', []))
    r.page_source(s, cfg.get('source'))


def build_blueprint(s, cfg, r):
    r.page_title(s, cfg['title'])
    lx, lw, cy = r.lx, r.lw, r.cy
    l2 = cfg['layer2']
    l1 = cfg['layer1']
    # Layer 2: AI Skills
    l2_y = cy + Inches(0.1)
    r.rrect(s, lx, l2_y, lw, Inches(1.9), 'blue_bg')
    r.rect(s, lx, l2_y, Inches(0.06), Inches(1.9), 'blue')
    r.txt(s, lx + Inches(0.2), l2_y + Inches(0.08), lw, Inches(0.3), l2['label'], sz=15, color='blue_dk', bold=True, font=r.pal['title_font'])
    skills = l2['skills']
    sw = (lw - Inches(0.4) - Inches(0.15) * (len(skills) - 1)) / len(skills)
    for i, sk in enumerate(skills):
        sx = lx + Inches(0.2) + i * (sw + Inches(0.15))
        r.rrect(s, sx, l2_y + Inches(0.5), sw, Inches(0.7), 'blue')
        r.txt(s, sx, l2_y + Inches(0.5), sw, Inches(0.7), sk, sz=13, color='white', bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    r.txt(s, lx + Inches(0.2), l2_y + Inches(1.35), lw - Inches(0.4), Inches(0.4),
          l2.get('footer', ''), sz=12, color='muted')
    # Connection
    r.txt(s, lx + lw / 2 - Inches(1.2), l2_y + Inches(1.95), Inches(2.4), Inches(0.35),
          cfg.get('connection', ''), sz=13, color='blue', bold=True, align=PP_ALIGN.CENTER)
    # Layer 1: MI + CRM
    l1_y = l2_y + Inches(2.35)
    r.rrect(s, lx, l1_y, lw, Inches(2.7), 'gray_bg')
    r.rect(s, lx, l1_y, Inches(0.06), Inches(2.7), 'ink')
    mi = l1['mi']
    mi_w = Inches(4.8)
    r.txt(s, lx + Inches(0.2), l1_y + Inches(0.08), mi_w, Inches(0.3), mi['label'], sz=15, color='ink', bold=True, font=r.pal['title_font'])
    mi_mods = mi['modules']
    mw = (mi_w - Inches(0.2) - Inches(0.12) * 2) / 3
    for i, mod in enumerate(mi_mods):
        col = i % 3; row = i // 3
        mx = lx + Inches(0.2) + col * (mw + Inches(0.12))
        my = l1_y + Inches(0.45) + row * Inches(0.75)
        r.rrect(s, mx, my, mw, Inches(0.6), 'white', 'line')
        r.txt(s, mx, my, mw, Inches(0.6), mod, sz=13, color='ink', bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    crm = l1['crm']
    mb_x = lx + mi_w + Inches(0.3)
    mb_w = lw - mi_w - Inches(0.5)
    r.txt(s, mb_x, l1_y + Inches(0.08), mb_w, Inches(0.3), crm['label'], sz=15, color='ink', bold=True, font=r.pal['title_font'])
    for i, end in enumerate(crm['ends']):
        ey = l1_y + Inches(0.45) + i * Inches(0.75)
        r.rrect(s, mb_x, ey, mb_w, Inches(0.6), 'white', 'line')
        r.txt(s, mb_x, ey, mb_w, Inches(0.6), end, sz=13, color='ink', bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    r.so_what(s, cfg.get('so_what', []))
    r.page_source(s, cfg.get('source'))


def build_domain_badge_cards(s, cfg, r):
    r.page_title(s, cfg['title'])
    lx, cy = r.lx, r.cy
    badges = cfg['badges']
    cols = cfg.get('columns', 3)
    for i, badge in enumerate(badges):
        col = i % cols; row = i // cols
        x = lx + col * Inches(2.5); y = cy + Inches(0.1) + row * Inches(1.5)
        c = badge.get('color', 'blue')
        r.rrect(s, x, y, Inches(2.3), Inches(1.2), 'blue_bg')
        r.rect(s, x, y, Inches(0.06), Inches(1.2), c)
        r.txt(s, x + Inches(0.2), y + Inches(0.12), Inches(2.0), Inches(0.35), badge['cn'],
              sz=18, color=c, bold=True, font=r.pal['title_font'])
        if badge.get('en'):
            r.txt(s, x + Inches(0.2), y + Inches(0.55), Inches(2.0), Inches(0.3), badge['en'], sz=12, color='muted')
    bb = cfg.get('bottom_bar')
    if bb:
        r.bottom_bar(s, bb['text'], bb.get('color', 'blue'))
    r.so_what(s, cfg.get('so_what', []))
    r.page_source(s, cfg.get('source'))


def build_domain_lifecycle(s, cfg, r):
    r.page_title(s, cfg['title'])
    lx, lw, cy = r.lx, r.lw, r.cy
    for ri, phase in enumerate(cfg['phases']):
        py = cy + Inches(0.1) + ri * Inches(1.6)
        color = phase.get('color', 'blue')
        r.rrect(s, lx, py, Inches(1.2), Inches(0.55), color)
        r.txt(s, lx, py, Inches(1.2), Inches(0.55), phase['name'], sz=14, color='white', bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        steps = phase['steps']
        n = len(steps)
        sw = (lw - Inches(1.4) - Inches(0.3) * (n - 1)) / n
        for si, step in enumerate(steps):
            sx = lx + Inches(1.4) + si * (sw + Inches(0.3))
            r.rrect(s, sx, py + Inches(0.03), sw, Inches(0.5), 'blue_bg', 'line')
            r.txt(s, sx, py + Inches(0.03), sw, Inches(0.5), step, sz=13, color='ink', bold=True,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            if si < n - 1:
                r.txt(s, sx + sw - Inches(0.05), py + Inches(0.08), Inches(0.3), Inches(0.4),
                      "→", sz=16, color=color, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    r.line(s, lx, cy + Inches(2.5), lw)
    for i, cap in enumerate(cfg.get('capabilities', [])):
        x = lx + i * Inches(2.55)
        r.txt(s, x, cy + Inches(2.7), Inches(2.3), Inches(0.3), cap['key'], sz=16, color='blue', bold=True, font=r.pal['title_font'])
        r.txt(s, x, cy + Inches(3.05), Inches(2.3), Inches(0.35), cap['val'], sz=13, color='body')
    bb = cfg.get('bottom_bar')
    if bb:
        r.bottom_bar(s, bb['text'], bb.get('color', 'blue'))
    r.so_what(s, cfg.get('so_what', []))
    r.page_source(s, cfg.get('source'))


def build_domain_big_number_flow(s, cfg, r):
    r.page_title(s, cfg['title'])
    lx, lw, cy = r.lx, r.lw, r.cy
    bn = cfg['big_number']
    r.big_number(s, lx, cy, Inches(3.5), bn['num'], bn['label'], bn.get('color', 'blue'))
    r.line(s, lx, cy + Inches(1.2), lw)
    flow = cfg['flow']
    for i, step in enumerate(flow):
        x = lx + i * Inches(1.55)
        r.txt(s, x, cy + Inches(1.5), Inches(1.4), Inches(0.3), str(i + 1), sz=13, color='muted', bold=True)
        r.txt(s, x, cy + Inches(1.75), Inches(1.4), Inches(0.5), step, sz=15, color='ink', bold=True)
        if i < len(flow) - 1:
            r.txt(s, x + Inches(1.25), cy + Inches(1.7), Inches(0.3), Inches(0.3), "→", sz=18, color='blue')
    if cfg.get('flow_description'):
        r.txt(s, lx, cy + Inches(2.5), lw, Inches(0.8), cfg['flow_description'], sz=14, color='body')
    bb = cfg.get('bottom_bar')
    if bb:
        r.bottom_bar(s, bb['text'], bb.get('color', 'blue'))
    r.so_what(s, cfg.get('so_what', []))
    r.page_source(s, cfg.get('source'))


def build_domain_comparison(s, cfg, r):
    r.page_title(s, cfg['title'])
    lx, cy = r.lx, r.cy
    before = cfg['before']; after = cfg['after']
    arrow = cfg.get('arrow', '→')
    # Before
    r.rrect(s, lx, cy, Inches(3.0), Inches(2.8), before.get('bg', 'gray_light'))
    r.txt(s, lx + Inches(0.2), cy + Inches(0.15), Inches(2.6), Inches(0.3), before['label'], sz=14, color='muted', bold=True, font=r.pal['title_font'])
    r.txt(s, lx + Inches(0.2), cy + Inches(0.5), Inches(2.6), Inches(0.8), before['value'], sz=48, color='muted', bold=True, font=r.pal['title_font'])
    r.txt(s, lx + Inches(0.2), cy + Inches(1.4), Inches(2.6), Inches(1.2), before['description'], sz=14, color='body', spacing=4)
    # Arrow
    r.txt(s, lx + Inches(3.1), cy + Inches(0.8), Inches(0.7), Inches(0.8),
          arrow, sz=48, color='blue', bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # After
    after_bg = after.get('bg', 'green_light')
    after_accent = after.get('accent', 'green')
    r.rrect(s, lx + Inches(3.9), cy, Inches(3.5), Inches(2.8), after_bg)
    r.rect(s, lx + Inches(3.9), cy, Inches(0.06), Inches(2.8), after_accent)
    r.txt(s, lx + Inches(4.1), cy + Inches(0.15), Inches(3.0), Inches(0.3), after['label'], sz=14, color=after_accent, bold=True, font=r.pal['title_font'])
    r.txt(s, lx + Inches(4.1), cy + Inches(0.5), Inches(3.0), Inches(0.8), after['value'], sz=48, color=after_accent, bold=True, font=r.pal['title_font'])
    r.txt(s, lx + Inches(4.1), cy + Inches(1.4), Inches(3.0), Inches(1.2), after['description'], sz=14, color='body', spacing=4)
    bb = cfg.get('bottom_bar')
    if bb:
        r.bottom_bar(s, bb['text'], bb.get('color', 'blue'))
    r.so_what(s, cfg.get('so_what', []))
    r.page_source(s, cfg.get('source'))


def build_domain_funnel(s, cfg, r):
    r.page_title(s, cfg['title'])
    lx, lw, cy = r.lx, r.lw, r.cy
    bn = cfg['big_number']
    r.big_number(s, lx, cy, Inches(3.0), bn['num'], bn['label'], bn.get('color', 'blue'))
    r.txt(s, lx + Inches(3.2), cy + Inches(0.05), Inches(4.5), Inches(0.35),
          cfg.get('funnel_title', ''), sz=16, color='blue', bold=True, font=r.pal['title_font'])
    items = cfg['funnel_items']
    for i, item in enumerate(items):
        y = cy + Inches(0.5) + i * Inches(0.42)
        w = Inches(4.5 - i * 0.5)
        r.rrect(s, lx + Inches(3.2), y, w, Inches(0.35), 'blue_bg')
        r.txt(s, lx + Inches(3.35), y + Inches(0.02), w, Inches(0.3), item, sz=14, color='ink', bold=True)
    r.line(s, lx, cy + Inches(2.8), lw)
    if cfg.get('summary'):
        r.txt(s, lx, cy + Inches(3.0), lw, Inches(0.35), cfg['summary'], sz=16, color='blue', bold=True, font=r.pal['title_font'])
    bb = cfg.get('bottom_bar')
    if bb:
        r.bottom_bar(s, bb['text'], bb.get('color', 'blue'))
    r.so_what(s, cfg.get('so_what', []))
    r.page_source(s, cfg.get('source'))


def build_domain_loop(s, cfg, r):
    r.page_title(s, cfg['title'])
    lx, lw, cy = r.lx, r.lw, r.cy
    bn = cfg['big_number']
    r.big_number(s, lx, cy, Inches(3.0), bn['num'], bn['label'], bn.get('color', 'blue'))
    loop = cfg['loop']
    for i, step in enumerate(loop):
        x = lx + Inches(3.2) + i * Inches(0.85)
        r.rrect(s, x, cy + Inches(0.1), Inches(0.75), Inches(0.55), 'blue')
        r.txt(s, x, cy + Inches(0.1), Inches(0.75), Inches(0.55), step, sz=14, color='white', bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < len(loop) - 1:
            r.txt(s, x + Inches(0.72), cy + Inches(0.15), Inches(0.15), Inches(0.4), "→", sz=16, color='blue', bold=True,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    r.line(s, lx, cy + Inches(1.3), lw)
    if cfg.get('secondary_title'):
        r.txt(s, lx, cy + Inches(1.5), lw, Inches(0.35), cfg['secondary_title'], sz=16, color='blue', bold=True, font=r.pal['title_font'])
    if cfg.get('description'):
        r.txt(s, lx, cy + Inches(1.9), lw, Inches(0.8), cfg['description'], sz=14, color='body')
    bb = cfg.get('bottom_bar')
    if bb:
        r.bottom_bar(s, bb['text'], bb.get('color', 'blue'))
    r.so_what(s, cfg.get('so_what', []))
    r.page_source(s, cfg.get('source'))


def build_product_ends(s, cfg, r):
    r.page_title(s, cfg['title'])
    lx, cy = r.lx, r.cy
    ends = cfg['ends']
    for i, end in enumerate(ends):
        x = lx + i * Inches(2.55)
        r.txt(s, x, cy, Inches(2.3), Inches(0.4), end['name'], sz=18, color='ink', bold=True, font=r.pal['title_font'])
        r.line(s, x, cy + Inches(0.45), Inches(0.6), 'blue', 2)
        r.txt(s, x, cy + Inches(0.6), Inches(2.3), Inches(1.5), end['desc'], sz=15, color='body', spacing=4)
    if cfg.get('reference_line'):
        r.txt(s, lx, cy + Inches(2.5), r.lw, Inches(0.4), cfg['reference_line'], sz=14, color='blue', bold=True)
    r.so_what(s, cfg.get('so_what', []))
    r.page_source(s, cfg.get('source'))


def build_chatbi_pipeline(s, cfg, r):
    r.page_title(s, cfg['title'])
    lx, cy = r.lx, r.cy
    step1 = cfg['step1']; step2 = cfg['step2']; step3 = cfg['step3']
    # Step 1
    s1_x = lx; s1_w = Inches(2.6)
    r.rrect(s, s1_x, cy + Inches(0.1), s1_w, Inches(4.2), 'blue_bg')
    r.rect(s, s1_x, cy + Inches(0.1), s1_w, Inches(0.5), step1.get('color', 'blue'))
    r.txt(s, s1_x, cy + Inches(0.1), s1_w, Inches(0.5), step1['title'], sz=14, color='white', bold=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for i, q in enumerate(step1['items']):
        qy = cy + Inches(0.8) + i * Inches(0.8)
        r.rrect(s, s1_x + Inches(0.15), qy, s1_w - Inches(0.3), Inches(0.6), 'white', 'line')
        r.txt(s, s1_x + Inches(0.25), qy, s1_w - Inches(0.5), Inches(0.6), q, sz=12, color='body', anchor=MSO_ANCHOR.MIDDLE)
    r.txt(s, s1_x + s1_w, cy + Inches(1.8), Inches(0.4), Inches(0.6), "→", sz=28, color='blue', bold=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Step 2
    s2_x = lx + Inches(3.0); s2_w = Inches(2.2)
    r.rrect(s, s2_x, cy + Inches(0.1), s2_w, Inches(4.2), 'gray_bg')
    r.rect(s, s2_x, cy + Inches(0.1), s2_w, Inches(0.5), step2.get('color', 'ink'))
    r.txt(s, s2_x, cy + Inches(0.1), s2_w, Inches(0.5), step2['title'], sz=14, color='white', bold=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for i, es in enumerate(step2['items']):
        ey = cy + Inches(0.8) + i * Inches(0.65)
        r.rrect(s, s2_x + Inches(0.15), ey, s2_w - Inches(0.3), Inches(0.45), 'white', 'line')
        r.txt(s, s2_x + Inches(0.15), ey, s2_w - Inches(0.3), Inches(0.45), es, sz=13, color='ink', bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    r.txt(s, s2_x + s2_w, cy + Inches(1.8), Inches(0.4), Inches(0.6), "→", sz=28, color='blue', bold=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Step 3
    s3_x = lx + Inches(5.5); s3_w = Inches(2.1)
    r.rrect(s, s3_x, cy + Inches(0.1), s3_w, Inches(4.2), 'green_light')
    r.rect(s, s3_x, cy + Inches(0.1), s3_w, Inches(0.5), step3.get('color', 'green'))
    r.txt(s, s3_x, cy + Inches(0.1), s3_w, Inches(0.5), step3['title'], sz=14, color='white', bold=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for i, item in enumerate(step3['items']):
        oy = cy + Inches(0.8) + i * Inches(0.8)
        r.txt(s, s3_x + Inches(0.1), oy, s3_w - Inches(0.2), Inches(0.25), item['type'], sz=13, color='green', bold=True)
        r.txt(s, s3_x + Inches(0.1), oy + Inches(0.25), s3_w - Inches(0.2), Inches(0.25), item['example'], sz=11, color='body')
    bb = cfg.get('bottom_bar')
    if bb:
        r.bottom_bar(s, bb['text'], bb.get('color', 'blue'))
    r.so_what(s, cfg.get('so_what', []))
    r.page_source(s, cfg.get('source'))


def build_ai_skill(s, cfg, r):
    title = (cfg.get('title_prefix', '') or '') + cfg['title']
    r.page_title(s, title)
    lx, lw, cy = r.lx, r.lw, r.cy
    flow = cfg['flow']
    for i, step in enumerate(flow):
        x = lx + Inches(0.3) + i * (Inches(7.0) / len(flow))
        w = Inches(7.0) / len(flow) - Inches(0.2)
        r.txt(s, x, cy + Inches(0.3), w, Inches(0.35), str(i + 1), sz=13, color='muted')
        r.txt(s, x, cy + Inches(0.6), w, Inches(0.4), step, sz=15, color='ink', bold=True)
        if i < len(flow) - 1:
            r.txt(s, x + w - Inches(0.05), cy + Inches(0.55), Inches(0.25), Inches(0.3), "→", sz=16, color='blue')
    r.line(s, lx, cy + Inches(1.3), lw)
    r.txt(s, lx, cy + Inches(1.6), lw, Inches(0.4), cfg['proof'], sz=16, color='blue', bold=True)
    r.so_what(s, cfg.get('so_what', []))
    r.page_source(s, cfg.get('source'))


def build_value_equation(s, cfg, r):
    r.page_title(s, cfg['title'])
    lx, lw, cy = r.lx, r.lw, r.cy
    items = cfg['items']
    operators = cfg.get('operators', ['−', '='])
    for i, item in enumerate(items):
        x = lx + i * Inches(2.55)
        color = item.get('color', 'blue')
        r.txt(s, x, cy + Inches(0.2), Inches(2.3), Inches(0.7), item['num'], sz=34, color=color, bold=True, font=r.pal['title_font'])
        r.txt(s, x, cy + Inches(0.95), Inches(2.3), Inches(0.35), item['label'], sz=15, color='ink', bold=True)
        r.txt(s, x, cy + Inches(1.35), Inches(2.3), Inches(1.0), item['desc'], sz=12, color='muted', spacing=2)
        if i < len(operators):
            r.txt(s, x + Inches(2.3), cy + Inches(0.4), Inches(0.25), Inches(0.5), operators[i], sz=28, color='ink', bold=True)
    r.line(s, lx, cy + Inches(2.8), lw)
    if cfg.get('comparison_old'):
        r.txt(s, lx, cy + Inches(3.1), lw, Inches(0.6), cfg['comparison_old'], sz=14, color='muted')
    if cfg.get('comparison_new'):
        r.txt(s, lx, cy + Inches(3.5), lw, Inches(0.6), cfg['comparison_new'], sz=14, color='blue', bold=True)
    r.so_what(s, cfg.get('so_what', []))
    r.page_source(s, cfg.get('source'))


def build_timeline(s, cfg, r):
    r.page_title(s, cfg['title'])
    lx, lw, cy = r.lx, r.lw, r.cy
    milestones = cfg['milestones']
    r.line(s, lx, cy + Inches(0.5), lw, 'blue', 2)
    for i, ms in enumerate(milestones):
        x = lx + Inches(0.3) + i * Inches(2.55)
        r.dot(s, x - Inches(0.05), cy + Inches(0.45), Inches(0.12), 'blue')
        r.txt(s, x, cy, Inches(2.3), Inches(0.3), ms['duration'], sz=14, color='blue', bold=True, font=r.pal['title_font'])
        r.txt(s, x, cy + Inches(0.7), Inches(2.3), Inches(0.35), ms['name'], sz=17, color='ink', bold=True, font=r.pal['title_font'])
        r.txt(s, x, cy + Inches(1.15), Inches(2.3), Inches(1.2), ms['desc'], sz=14, color='body', spacing=3)
    r.so_what(s, cfg.get('so_what', []))
    r.page_source(s, cfg.get('source'))


def build_timeline_detail(s, cfg, r):
    r.page_title(s, cfg['title'])
    lx, lw, cy = r.lx, r.lw, r.cy
    milestones = cfg['milestones']
    r.line(s, lx, cy + Inches(0.5), lw, 'blue', 2)
    for i, ms in enumerate(milestones):
        x = lx + Inches(0.3) + i * Inches(2.55)
        r.dot(s, x - Inches(0.05), cy + Inches(0.45), Inches(0.12), 'blue')
        r.txt(s, x, cy, Inches(2.3), Inches(0.3), ms['duration'], sz=14, color='blue', bold=True, font=r.pal['title_font'])
        r.txt(s, x, cy + Inches(0.7), Inches(2.3), Inches(0.35), ms['name'], sz=17, color='ink', bold=True, font=r.pal['title_font'])
        for j, item in enumerate(ms['items']):
            y = cy + Inches(1.2) + j * Inches(0.5)
            r.dot(s, x, y + Inches(0.08), Inches(0.07))
            r.txt(s, x + Inches(0.18), y, Inches(2.1), Inches(0.4), item, sz=13, color='body')
    r.so_what(s, cfg.get('so_what', []))
    r.page_source(s, cfg.get('source'))


def build_flywheel(s, cfg, r):
    r.page_title(s, cfg['title'])
    lx, cy = r.lx, r.cy
    cycle = cfg['cycle']
    cx, cy_center = lx + Inches(3.3), cy + Inches(2.0)
    radius = Inches(1.8)
    for i, step in enumerate(cycle):
        angle = i * (360 / len(cycle))
        x = cx + Inches(radius.inches * math.cos(math.radians(angle - 90))) - Inches(0.7)
        y = cy_center + Inches(radius.inches * math.sin(math.radians(angle - 90))) - Inches(0.15)
        r.dot(s, x + Inches(0.55), y, Inches(0.12), 'blue')
        r.txt(s, x, y + Inches(0.2), Inches(1.4), Inches(0.3), step, sz=14, color='ink', bold=True, align=PP_ALIGN.CENTER)
    center_label = cfg.get('center_label', '飞轮\n效应')
    r.txt(s, lx + Inches(2.3), cy + Inches(1.7), Inches(2), Inches(0.5), center_label,
          sz=18, color='blue', bold=True, font=r.pal['title_font'], align=PP_ALIGN.CENTER)
    if cfg.get('footer'):
        r.txt(s, lx, cy + Inches(4.3), r.lw, Inches(0.4), cfg['footer'], sz=15, color='ink', bold=True)
    r.so_what(s, cfg.get('so_what', []))
    r.page_source(s, cfg.get('source'))


def build_pricing_cards(s, cfg, r):
    r.page_title(s, cfg['title'])
    lx, cy = r.lx, r.cy
    cards = cfg['cards']
    for i, card in enumerate(cards):
        x = lx + i * Inches(2.55)
        color = card.get('color', 'ink')
        r.txt(s, x, cy, Inches(2.3), Inches(0.6), card['price'], sz=28, color=color, bold=True, font=r.pal['title_font'])
        r.line(s, x, cy + Inches(0.7), Inches(0.6), color, 2)
        r.txt(s, x, cy + Inches(0.85), Inches(2.3), Inches(0.35), card['label'], sz=16, color='ink', bold=True, font=r.pal['title_font'])
        r.txt(s, x, cy + Inches(1.3), Inches(2.3), Inches(1.5), card['desc'], sz=14, color='body', spacing=4)
    if cfg.get('footer'):
        r.txt(s, lx, cy + Inches(3.5), r.lw, Inches(0.5), cfg['footer'], sz=16, color='blue', bold=True, align=PP_ALIGN.CENTER)
    r.so_what(s, cfg.get('so_what', []))
    r.page_source(s, cfg.get('source'))


def build_closing_pillars(s, cfg, r):
    r.page_title(s, cfg['title'])
    pillars = cfg['pillars']
    p_total = Inches(1.8) * len(pillars) + Inches(0.4) * (len(pillars) - 1)
    p_start = (Inches(13.333) - p_total) / 2
    cy = r.cy
    for i, pillar in enumerate(pillars):
        x = p_start + i * Inches(2.2)
        r.txt(s, x, cy + Inches(0.5), Inches(1.8), Inches(0.7), pillar['keyword'],
              sz=36, color='blue', bold=True, font=r.pal['title_font'], align=PP_ALIGN.CENTER)
        r.txt(s, x, cy + Inches(1.3), Inches(1.8), Inches(0.35), pillar['label'], sz=15, color='ink', bold=True, align=PP_ALIGN.CENTER)
        r.line(s, x + Inches(0.4), cy + Inches(1.8), Inches(1.0))
        r.txt(s, x, cy + Inches(2.0), Inches(1.8), Inches(1.0), pillar['proof'], sz=13, color='muted', align=PP_ALIGN.CENTER, spacing=3)
    if cfg.get('footer'):
        r.txt(s, Inches(1.5), cy + Inches(3.5), Inches(10.3), Inches(0.8), cfg['footer'],
              sz=18, color='ink', bold=True, align=PP_ALIGN.CENTER, spacing=6)
    r.page_source(s, cfg.get('source'))


def build_product_overview(s, cfg, r):
    r.page_title(s, cfg['title'])
    tx, tw, cy = r.tx, r.tw, r.cy
    products = cfg['products']
    col_w = (tw - Inches(0.4)) / len(products)
    for i, prod in enumerate(products):
        x = tx + i * (col_w + Inches(0.2))
        color = prod.get('color', 'blue')
        r.rrect(s, x, cy, col_w, Inches(0.7), color)
        r.txt(s, x, cy, col_w, Inches(0.7), prod['name'], sz=20, color='white', bold=True, font=r.pal['title_font'],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        r.txt(s, x, cy + Inches(0.8), col_w, Inches(0.3), prod['tagline'], sz=14, color='muted', align=PP_ALIGN.CENTER)
        features = prod['features']
        if len(features) > 6:
            half = (len(features) + 1) // 2
            sub_w = (col_w - Inches(0.3)) / 2
            for j, feat in enumerate(features):
                col_j = j // half; row_j = j % half
                fx = x + Inches(0.1) + col_j * (sub_w + Inches(0.1))
                fy = cy + Inches(1.3) + row_j * Inches(0.5)
                r.rrect(s, fx, fy, sub_w, Inches(0.4), 'blue_bg')
                r.txt(s, fx, fy, sub_w, Inches(0.4), feat, sz=12, color='ink', bold=True,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        else:
            for j, feat in enumerate(features):
                fy = cy + Inches(1.3) + j * Inches(0.5)
                r.rrect(s, x + Inches(0.1), fy, col_w - Inches(0.2), Inches(0.4), 'blue_bg')
                r.txt(s, x + Inches(0.1), fy, col_w - Inches(0.2), Inches(0.4), feat, sz=13, color='ink', bold=True,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if cfg.get('footer'):
        r.txt(s, tx, cy + Inches(5.0), tw, Inches(0.4), cfg['footer'], sz=18, color='blue', bold=True, align=PP_ALIGN.CENTER)
    r.page_source(s, cfg.get('source', '蓝联科技 · 产品全景'))


# =============================================================================
# Slide Type Registry
# =============================================================================

SLIDE_BUILDERS = {
    'requirement_projects':     build_requirement_projects,
    'requirement_roles':        build_requirement_roles,
    'requirement_pain_detail':  build_requirement_pain_detail,
    'requirement_matrix':       build_requirement_matrix,
    'blueprint':                build_blueprint,
    'domain_badge_cards':       build_domain_badge_cards,
    'domain_lifecycle':         build_domain_lifecycle,
    'domain_big_number_flow':   build_domain_big_number_flow,
    'domain_comparison':        build_domain_comparison,
    'domain_funnel':            build_domain_funnel,
    'domain_loop':              build_domain_loop,
    'product_ends':             build_product_ends,
    'chatbi_pipeline':          build_chatbi_pipeline,
    'ai_skill':                 build_ai_skill,
    'value_equation':           build_value_equation,
    'timeline':                 build_timeline,
    'timeline_detail':          build_timeline_detail,
    'flywheel':                 build_flywheel,
    'pricing_cards':            build_pricing_cards,
    'closing_pillars':          build_closing_pillars,
    'product_overview':         build_product_overview,
}


# =============================================================================
# TOC helpers
# =============================================================================

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _grad(run, c1, c2):
    rPr = run._r.get_or_add_rPr()
    for tag in ("gradFill", "solidFill", "noFill"):
        for e in rPr.findall(f"{{{A_NS}}}{tag}"):
            rPr.remove(e)
    g = etree.fromstring(
        f'<a:gradFill xmlns:a="{A_NS}"><a:gsLst>'
        f'<a:gs pos="44000"><a:srgbClr val="{c1}"/></a:gs>'
        f'<a:gs pos="100000"><a:srgbClr val="{c2}"/></a:gs>'
        f'</a:gsLst><a:lin ang="5400000" scaled="0"/></a:gradFill>'
    )
    rPr.insert(0, g)


def set_toc(slide, sections, active):
    shapes = sorted(
        [s for s in slide.shapes if getattr(s, "has_text_frame", False)],
        key=lambda s: s.top or 0
    )
    if len(shapes) < 3:
        return
    n = min(len(sections), 4)
    for i in range(min(len(shapes), n)):
        tf = shapes[i].text_frame; tf.clear()
        p = tf.paragraphs[0]; p.text = f"{i+1}. {sections[i]}"; p.alignment = PP_ALIGN.LEFT
        if not p.runs:
            p.add_run()
        r = p.runs[0]; r.font.size = Pt(32); r.font.bold = True; r.font.name = "Microsoft YaHei"
        if (i + 1) == active:
            _grad(r, "2C5DE6", "6F92F3")
        else:
            r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


# =============================================================================
# Main Compiler
# =============================================================================

def compile_v5(config_path):
    """Compile a V5 PPT from a YAML config file."""

    # ── Load config ──
    with open(config_path, 'r', encoding='utf-8') as f:
        cfg = yaml.safe_load(f)

    meta = cfg.get('meta', {})
    paths = cfg.get('paths', {})
    r = Renderer(cfg.get('palette'), cfg.get('geometry'))

    base_path = _expand_path(paths.get('base_ppt'))
    output_path = _expand_path(paths.get('output'))
    if not base_path or not output_path:
        print("错误: paths.base_ppt 和 paths.output 必填")
        sys.exit(1)

    print(f"V5 领域驱动方案编译器")
    print(f"  客户: {meta.get('customer_name', '?')}")
    print(f"  Base: {base_path}")
    print(f"  Output: {output_path}")
    print()

    # ── Load base PPT ──
    prs = Presentation(base_path)
    nb = len(prs.slides)
    print(f"  Base PPT: {nb} slides")
    layout = prs.slide_layouts[2]

    # ── Fix P3 text overflow (optional) ──
    p3_fix = cfg.get('p3_fix', {})
    if p3_fix.get('enabled'):
        match = p3_fix.get('match_text', '')
        size = p3_fix.get('font_size', 11)
        wrap = p3_fix.get('word_wrap', True)
        for sh in prs.slides[2].shapes:
            if getattr(sh, "has_text_frame", False) and match in sh.text_frame.text and len(sh.text_frame.text) > 100:
                sh.text_frame.word_wrap = wrap
                for p in sh.text_frame.paragraphs:
                    for run in p.runs:
                        if not run.font.size:
                            run.font.size = Pt(size)
                break
        print(f"  ✎ P3 text overflow fixed")

    # ── Fix cover ──
    cover_cfg = cfg.get('cover', {})
    if cover_cfg:
        match_text = cover_cfg.get('match_text', '')
        new_title = cover_cfg.get('new_title', '')
        title_size = cover_cfg.get('title_size', 40)
        for sh in prs.slides[0].shapes:
            if getattr(sh, "has_text_frame", False) and match_text in sh.text_frame.text:
                sh.text_frame.paragraphs[0].text = new_title
                for run in sh.text_frame.paragraphs[0].runs:
                    run.font.size = Pt(title_size)
                    run.font.color.rgb = r.color('white')
                    run.font.bold = True
                break
        if cover_cfg.get('subtitle'):
            r.txt(prs.slides[0], Inches(2.5), Inches(cover_cfg.get('subtitle_y', 5.2)),
                  Inches(8.3), Inches(0.6), cover_cfg['subtitle'], sz=18, color='blue', align=PP_ALIGN.CENTER)
        if meta.get('date'):
            r.txt(prs.slides[0], Inches(2.5), Inches(cover_cfg.get('date_y', 5.9)),
                  Inches(8.3), Inches(0.4), f"方案日期：{meta['date']}", sz=13, color='muted', align=PP_ALIGN.CENTER)
        print(f"  ✎ Cover updated")

    # ── Build custom slides ──
    slides_cfg = cfg.get('slides', [])
    if not isinstance(slides_cfg, list):
        slides_cfg = []
    custom_count = 0
    for slide_cfg in slides_cfg:
        stype = slide_cfg.get('type')
        builder = SLIDE_BUILDERS.get(stype)
        if not builder:
            print(f"  ⚠ Unknown slide type: {stype}, skipping")
            continue
        slide = prs.slides.add_slide(layout)
        builder(slide, slide_cfg, r)
        custom_count += 1
    print(f"  + {custom_count} custom slides")

    # ── Build product overview (optional, replaces a base slide) ──
    prod_idx = None
    po_cfg = cfg.get('product_overview')
    if po_cfg:
        slide = prs.slides.add_slide(layout)
        build_product_overview(slide, po_cfg, r)
        prod_idx = len(prs.slides) - 1
        print(f"  + 1 product overview slide (replaces base index {po_cfg.get('replace_base_index', '?')})")

    # ── Fix TOC dividers ──
    toc_cfg = cfg.get('toc', {})
    if toc_cfg:
        sections = toc_cfg.get('sections', [])
        divider_indices = toc_cfg.get('divider_indices', [])
        for i, idx in enumerate(divider_indices):
            if idx < nb:
                set_toc(prs.slides[idx], sections, i + 1)
        print(f"  ✎ {len(divider_indices)} TOC dividers corrected")

    # ── Reorder ──
    reorder_cfg = cfg.get('reorder', {})
    if reorder_cfg:
        _do_reorder(prs, reorder_cfg, nb, prod_idx)
        print(f"  ✎ Slides reordered")

    # ── Save ──
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    sz = os.path.getsize(output_path) // 1024
    print(f"\n  Output: {output_path}")
    print(f"  Total: {len(prs.slides)} slides ({sz}KB)")
    print(f"\n✓ Done")


def _do_reorder(prs, cfg, nb, prod_idx):
    """Reorder slides: cover + toc + company_intro + [divider + custom_section]* + closing."""
    cover_idx = cfg.get('cover_index', 0)
    toc_idx = cfg.get('toc_index', 1)
    company_range = cfg.get('company_intro_range', [])
    closing_idx = cfg.get('closing_index', 61)
    dividers = cfg.get('section_dividers', [])
    section_counts = cfg.get('section_slide_counts', [])
    po_replace = cfg.get('product_overview_replace_index')

    # Build company intro slide list
    s1 = list(company_range)
    if prod_idx is not None and po_replace is not None and po_replace < len(s1):
        s1[po_replace] = prod_idx

    # Build result order
    result = [cover_idx, toc_idx] + s1
    off = nb  # custom slides start after base slides
    for i, count in enumerate(section_counts):
        if i < len(dividers):
            result.append(dividers[i])
        result += list(range(off, off + count))
        off += count
    result.append(closing_idx)

    # Execute reorder via XML manipulation
    xml = prs.slides._sldIdLst
    elems = list(xml)
    total = len(elems)

    ordered = [elems[i] for i in result if i < total]
    ordered_set = set(i for i in result if i < total)
    leftovers = [elems[i] for i in range(total) if i not in ordered_set]

    for e in elems:
        xml.remove(e)
    for e in ordered + leftovers:
        xml.append(e)
    # Drop leftover slides (unused base slides)
    for e in leftovers:
        rId = e.get(qn("r:id"))
        xml.remove(e)
        if rId:
            prs.part.drop_rel(rId)


# =============================================================================
# Entry point
# =============================================================================

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python compile_v5_domain.py <config.yaml>")
        print("配置规范: domain-slide-config-schema.yaml")
        sys.exit(1)
    config_path = sys.argv[1]
    if not os.path.exists(config_path):
        print(f"错误: 配置文件不存在: {config_path}")
        sys.exit(1)
    compile_v5(config_path)
