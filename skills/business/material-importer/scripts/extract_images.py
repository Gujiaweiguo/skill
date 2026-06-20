"""从 PPT/Word/Excel 中提取图片并修正 Markdown 图片路径。支持 --json 输出供 Agent 程序化读取。"""

import json
import os
import re
import sys
import zipfile
from pathlib import Path


def extract_images_from_zip(zip_path: str, media_dir: str) -> dict[str, str]:
    """从 ZIP 中提取图片，返回 原名→新名 的映射。"""
    os.makedirs(media_dir, exist_ok=True)
    name_map = {}

    try:
        with zipfile.ZipFile(zip_path, "r") as z:
            # 判断是 PPTX 还是 DOCX
            ppt_media = [n for n in z.namelist() if n.startswith("ppt/media/") and "/" != n]
            doc_media = [n for n in z.namelist() if n.startswith("word/media/") and "/" != n]
            media_files = ppt_media or doc_media

            for name in media_files:
                try:
                    data = z.read(name)
                except Exception:
                    continue  # 跳过坏文件

                basename = os.path.basename(name)  # e.g. image12.jpeg
                base_noext, ext = os.path.splitext(basename)  # image12, .jpeg
                dest = os.path.join(media_dir, basename)

                if not os.path.exists(dest):
                    with open(dest, "wb") as f:
                        f.write(data)

                # 记录原名：PPT 中引用关系 id 映射，简化处理
                # 用序号作为 key 的备选
                seq = base_noext.replace("image", "")  # "12"
                name_map[seq] = basename

    except Exception as e:
        print(f"  ⚠️ 提取失败: {e}")

    return name_map


def extract_images_pptx(pptx_path: str, media_dir: str) -> dict[str, str]:
    """使用 python-pptx 提取图片并获取关系映射。"""
    from pptx import Presentation
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    os.makedirs(media_dir, exist_ok=True)
    rel_map = {}  # rId → 文件名

    try:
        prs = Presentation(pptx_path)

        # 遍历所有幻灯片和形状，收集图片关系
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.shape_type == 13:  # PICTURE
                    try:
                        image = shape.image
                        # 图片原始文件名（来自 PPT 内部关系）
                        blob = image.blob
                        content_type = image.content_type
                        ext = content_type.split("/")[-1] if "/" in content_type else "png"
                        if ext == "jpeg":
                            ext = "jpg"
                        if ext == "png" or ext == "jpg":
                            pass
                        filename = f"image_s{slide_num}_{shape.shape_id}.{ext}"
                        dest = os.path.join(media_dir, filename)
                        if not os.path.exists(dest):
                            with open(dest, "wb") as f:
                                f.write(blob)
                        # 同时索引图片的原始关系名
                        rel_id = image._rId if hasattr(image, '_rId') else ""
                        if rel_id:
                            rel_map[rel_id] = filename
                    except Exception:
                        continue

        # 尝试从 ZIP 直接读取关系 XML 建立 rId → 文件名映射
        try:
            with zipfile.ZipFile(pptx_path, "r") as z:
                # 读取 ppt/_rels/presentation.xml.rels 获取图片关系
                import xml.etree.ElementTree as ET
                rels_xml = z.read("ppt/_rels/presentation.xml.rels")
                root = ET.fromstring(rels_xml)
                ns = {"r": "http://schemas.openxmlformats.org/package/2006/relationships"}

                # 从每个 slide 的 rels 中读取
                slides = [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
                for slide_path in slides:
                    rels_path = slide_path.replace("/slide", "/_rels/slide").replace(".xml", ".xml.rels")
                    if rels_path in z.namelist():
                        try:
                            rels = ET.fromstring(z.read(rels_path))
                            for child in rels:
                                r_id = child.attrib.get("Id", "")
                                target = child.attrib.get("Target", "")
                                # 目标可能是 ../media/image12.jpeg
                                if "media/" in target:
                                    media_name = target.split("/")[-1]
                                    rel_map[r_id] = media_name
                        except Exception:
                            continue
        except Exception:
            pass

    except Exception as e:
        print(f"  ⚠️ python-pptx 处理失败: {e}")

    return rel_map


def extract_images_from_xlsx(xlsx_path: str, media_dir: str) -> None:
    """从 xlsx 提取图片（xl/media/ + 各 sheet 嵌入图）。"""
    import zipfile
    os.makedirs(media_dir, exist_ok=True)

    # 1. ZIP 直接提取 xl/media/
    try:
        with zipfile.ZipFile(xlsx_path, "r") as z:
            for name in z.namelist():
                if name.startswith("xl/media/") and "/" != name:
                    basename = os.path.basename(name)
                    dest = os.path.join(media_dir, basename)
                    if not os.path.exists(dest):
                        data = z.read(name)
                        with open(dest, "wb") as f:
                            f.write(data)
    except Exception as e:
        print(f"  ⚠️ xl/media/ 提取失败: {e}")

    # 2. openpyxl sheet._images 补充提取（部分 xlsx 的图片挂在 sheet 上）
    try:
        from openpyxl import load_workbook
        wb = load_workbook(xlsx_path)
        for sn in wb.sheetnames:
            if sn == "Sheet1":
                continue
            ws = wb[sn]
            idx = 1
            for img in getattr(ws, "_images", []):
                data = img.ref.read()
                fname = f"{sn}_contract_{idx}.png"
                dest = os.path.join(media_dir, fname)
                if not os.path.exists(dest):
                    with open(dest, "wb") as f:
                        f.write(data)
                idx += 1
    except Exception as e:
        print(f"  ⚠️ openpyxl sheet 图片提取: {e}")


def fix_markdown_paths(md_path: str, media_dir: str, rel_map: dict[str, str], stem: str) -> int:
    """修正 Markdown 中的图片路径，返回修正数量。"""
    if not os.path.exists(md_path):
        return 0

    with open(md_path, "r", encoding="utf-8") as f:
        content = f.read()

    original = content
    fix_count = 0

    # 匹配 ![alt](文件名)
    def replace_image(match):
        nonlocal fix_count
        alt = match.group(1)
        ref = match.group(2).strip()
        ref_basename = os.path.basename(ref)

        # 策略 1: 直接查找
        dest_path = os.path.join(media_dir, ref_basename)
        if os.path.exists(dest_path):
            fix_count += 1
            return f"![{alt}](raw/{stem}_media/{ref_basename})"

        # 策略 2: 去掉路径前缀再查找
        bare_ref = os.path.basename(ref)
        dest_path = os.path.join(media_dir, bare_ref)
        if os.path.exists(dest_path):
            fix_count += 1
            return f"![{alt}](raw/{stem}_media/{bare_ref})"

        # 策略 3: 遍历 media 目录，按名称相似性匹配
        ref_noext = os.path.splitext(bare_ref)[0]
        for fname in os.listdir(media_dir):
            fname_noext = os.path.splitext(fname)[0]
            # 去除数字前缀差异匹配
            ref_digits = re.sub(r"\D", "", ref_noext)
            fname_digits = re.sub(r"\D", "", fname_noext)
            if ref_digits and ref_digits == fname_digits:
                fix_count += 1
                return f"![{alt}](raw/{stem}_media/{fname})"

        return match.group(0)

    content = re.sub(r"!\[([^\]]*)\]\(([^)]+)\)", replace_image, content)

    if content != original:
        with open(md_path, "w", encoding="utf-8") as f:
            f.write(content)
        print(f"  ✅ 图片路径已修正 {fix_count} 处: {stem}")
    else:
        print(f"  ⚠️ 未找到可匹配的图片: {stem}")

    return fix_count


def process_file(file_path: str, raw_dir: str) -> None:
    """处理单个文件：提取图片 + 修正 Markdown 路径。"""
    ext = Path(file_path).suffix.lower()
    stem = Path(file_path).stem
    media_dir = os.path.join(raw_dir, f"{stem}_media")
    md_path = os.path.join(raw_dir, f"{Path(file_path).name}.md")

    if ext == ".pptx":
        extracted = extract_images_pptx(file_path, media_dir)
        total_images = len([f for f in os.listdir(media_dir) if os.path.isfile(os.path.join(media_dir, f))]) if os.path.exists(media_dir) else 0
        print(f"  📷 PPT 图片: {total_images} 张 -> {stem}_media/")
        fix_markdown_paths(md_path, media_dir, extracted, stem)

    elif ext == ".docx":
        name_map = extract_images_from_zip(file_path, media_dir)
        total_images = len([f for f in os.listdir(media_dir) if os.path.isfile(os.path.join(media_dir, f))]) if os.path.exists(media_dir) else 0
        print(f"  📷 Word 图片: {total_images} 张 -> {stem}_media/")
        fix_markdown_paths(md_path, media_dir, name_map, stem)

    elif ext == ".xlsx":
        extract_images_from_xlsx(file_path, media_dir)
        total_images = len([f for f in os.listdir(media_dir) if os.path.isfile(os.path.join(media_dir, f))]) if os.path.exists(media_dir) else 0
        print(f"  📷 Excel 图片: {total_images} 张 -> {stem}_media/")

    else:
        print(f"  ⏭️ 跳过 (不支持格式): {file_path}")


def main():
    import argparse
    parser = argparse.ArgumentParser(description="从 PPT/Word/Excel 提取图片")
    parser.add_argument("target", help="文件或目录路径")
    parser.add_argument("--json", action="store_true", help="JSON 输出")
    parser.add_argument("--raw-dir", help="raw 输出目录（默认: target 同级的 raw/）")
    parsed = parser.parse_args()

    use_json = parsed.json
    target = parsed.target.rstrip("/\\")
    raw_dir = parsed.raw_dir or (os.path.join(target, "raw") if os.path.isdir(target) else os.path.dirname(target))

    results = []

    if os.path.isdir(target):
        for f in sorted(os.listdir(target)):
            if f.startswith("~$") or not f.endswith((".pptx", ".docx", ".xlsx")):
                continue
            file_path = os.path.join(target, f)
            process_file(file_path, raw_dir)
            stem = Path(file_path).stem
            media_dir = os.path.join(raw_dir, f"{stem}_media")
            img_count = 0
            if os.path.exists(media_dir):
                img_count = len([x for x in os.listdir(media_dir) if os.path.isfile(os.path.join(media_dir, x))])
            results.append({"file": f, "stem": stem, "images": img_count})
    else:
        process_file(target, raw_dir)
        stem = Path(target).stem
        media_dir = os.path.join(raw_dir, f"{stem}_media")
        img_count = 0
        if os.path.exists(media_dir):
            img_count = len([x for x in os.listdir(media_dir) if os.path.isfile(os.path.join(media_dir, x))])
        results.append({"file": os.path.basename(target), "stem": stem, "images": img_count})

    if use_json:
        report = {"tool": "extract_images", "files": results}
        print(json.dumps(report, ensure_ascii=False, indent=2))
    else:
        for r in results:
            print(f"  📄 {r['file']}: {r['images']} 张图片")
        print("\n✅ 图片提取完成")


if __name__ == "__main__":
    main()