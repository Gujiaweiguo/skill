"""
文档读取器 — 格式检测 + 多格式读取 → RawDocument

支持的格式:
  .docx  → python-docx（含图片提取）
  .pptx  → python-pptx（含图片提取）
  .xlsx  → openpyxl
  .pdf   → LibreOffice → .docx → python-docx（仅文字型 PDF）
  .doc   → ❌ 请用户用 WPS/Office 另存为 .docx

设计原则：
  - 不依赖 markitdown
  - 不处理 OCR
  - 图片提取到 assets/ 目录供内容包引用
"""

import os
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from .models import ImageAsset, RawDocument, TableData

# ── 格式检测 ──────────────────────────────────────────────

SUPPORTED_FORMATS = {
    ".docx": "docx",
    ".pptx": "pptx",
    ".xlsx": "xlsx",
    ".xls": "xls",
    ".pdf": "pdf",
}


def detect_format(file_path: str) -> str:
    """检测文件格式（基于扩展名）"""
    ext = Path(file_path).suffix.lower()
    fmt = SUPPORTED_FORMATS.get(ext)
    if not fmt:
        raise ValueError(
            f"不支持的文件格式: {ext}\n"
            f"支持: {', '.join(SUPPORTED_FORMATS)}\n"
            f"提示: .doc 请用 WPS/Office 另存为 .docx 后重试"
        )
    return fmt


# ── 主入口 ────────────────────────────────────────────────


def read_document(file_path: str, assets_dir: str = "", extract_images: bool = False) -> RawDocument:
    """读取文档 → RawDocument

    Args:
        file_path: 文件路径
        assets_dir: 图片提取目录（默认: 同目录下 {文件名}_images/）
        extract_images: 是否提取图片（默认 False，招标文件读取不需要）

    Returns:
        RawDocument 对象
    """
    file_path = os.path.abspath(file_path)
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")

    fmt = detect_format(file_path)

    # 自动确定 assets 目录
    if not assets_dir:
        stem = Path(file_path).stem
        assets_dir = os.path.join(os.path.dirname(file_path), f"{stem}_images")

    doc = RawDocument(file_path=file_path, file_type=fmt)

    if fmt == "docx":
        _read_docx(file_path, doc, assets_dir, extract_images=extract_images)
    elif fmt == "pptx":
        _read_pptx(file_path, doc, assets_dir, extract_images=extract_images)
    elif fmt == "xlsx":
        _read_xlsx(file_path, doc)
    elif fmt == "xls":
        _read_xls(file_path, doc)
    elif fmt == "pdf":
        _read_pdf(file_path, doc, assets_dir, extract_images=extract_images)

    if not doc.title:
        doc.title = _guess_title(doc.text_content, file_path)

    return doc


# ── .docx 读取 ────────────────────────────────────────────


def _read_docx(file_path: str, doc: RawDocument, assets_dir: str, extract_images: bool = False):
    """读取 .docx：默认用 zipfile+xml（快且稳），需要图片时用 python-docx"""
    if extract_images:
        _read_docx_fast(file_path, doc, assets_dir, extract_images)
    else:
        _read_docx_fallback(file_path, doc)


def _read_docx_fast(file_path: str, doc: RawDocument, assets_dir: str, extract_images: bool = False):
    """标准 python-docx 读取（速度快）"""
    from docx import Document

    d = Document(file_path)

    # 段落
    paragraphs = []
    for p in d.paragraphs:
        text = p.text.strip()
        if not text:
            paragraphs.append("")
            continue
        style_name = (p.style.name or "") if p.style else ""
        if "Heading" in style_name or "标题" in style_name:
            try:
                level = style_name.replace("Heading ", "").replace("标题 ", "").strip()
                prefix = "#" * int(level)
            except (ValueError, IndexError):
                prefix = "##"
            paragraphs.append(f"{prefix} {text}")
        else:
            paragraphs.append(text)

    # 表格
    for table in d.tables:
        td = TableData()
        for i, row in enumerate(table.rows):
            cells = [cell.text.strip() for cell in row.cells]
            if i == 0:
                td.headers = cells
            else:
                td.rows.append(cells)
        if td.headers or td.rows:
            doc.tables.append(td)

    # 图片（可选）
    if extract_images:
        _extract_docx_images(d, assets_dir, doc)

    doc.text_content = "\n".join(paragraphs)


def _read_docx_fallback(file_path: str, doc: RawDocument):
    """降级方案：直接用 zipfile + xml 解析 word/document.xml"""
    import zipfile
    import xml.etree.ElementTree as ET

    NS = {
        "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
        "wp": "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing",
    }

    paragraphs = []
    tables = []

    with zipfile.ZipFile(file_path, "r") as zf:
        if "word/document.xml" not in zf.namelist():
            doc.text_content = "[无法读取]: 无效的 docx 文件（缺少 word/document.xml）"
            return

        with zf.open("word/document.xml") as f:
            tree = ET.parse(f)
            root = tree.getroot()

            # 遍历所有段落和表格
            body = root.find(".//w:body", NS)
            if body is None:
                doc.text_content = "[无法读取]: 文档结构异常"
                return

            for child in body:
                tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

                if tag == "p":
                    # 段落
                    texts = []
                    for t in child.iter("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t"):
                        if t.text:
                            texts.append(t.text)
                    text = "".join(texts).strip()
                    if text:
                        # 简单判断标题：看 pPr 里的 pStyle
                        pPr = child.find("w:pPr", NS)
                        is_heading = False
                        if pPr is not None:
                            pStyle = pPr.find("w:pStyle", NS)
                            if pStyle is not None:
                                style_val = pStyle.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val", "")
                                if "Heading" in style_val or "Title" in style_val:
                                    is_heading = True
                                    level = style_val.replace("Heading", "").strip() or "1"
                                    try:
                                        prefix = "#" * int(level)
                                    except ValueError:
                                        prefix = "##"
                                    paragraphs.append(f"{prefix} {text}")
                        if not is_heading:
                            paragraphs.append(text)
                    else:
                        paragraphs.append("")

                elif tag == "tbl":
                    # 表格
                    td = TableData()
                    for row_idx, tr in enumerate(child.iter("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}tr")):
                        cells = []
                        for tc in tr.iter("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}tc"):
                            cell_texts = []
                            for t in tc.iter("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t"):
                                if t.text:
                                    cell_texts.append(t.text)
                            cells.append("".join(cell_texts).strip())
                        if row_idx == 0:
                            td.headers = cells
                        else:
                            td.rows.append(cells)
                    if td.headers or td.rows:
                        tables.append(td)

    doc.text_content = "\n".join(paragraphs)
    doc.tables = tables


def _extract_docx_images(document, assets_dir: str, doc: RawDocument):
    """从 docx 中提取图片"""
    from docx.oxml.ns import qn

    os.makedirs(assets_dir, exist_ok=True)
    img_count = 0

    for para_idx, para in enumerate(document.paragraphs):
        for run in para.runs:
            drawing_elements = run._element.findall(qn("w:drawing"))
            for drawing in drawing_elements:
                blip = drawing.findall(
                    ".//" + qn("a:blip")
                )
                for b in blip:
                    embed_id = b.get(qn("r:embed"))
                    if not embed_id:
                        continue
                    try:
                        image_part = document.part.related_parts[embed_id]
                        ext = _image_ext_from_content_type(image_part.content_type)
                        img_filename = f"image_{img_count:03d}.{ext}"
                        img_path = os.path.join(assets_dir, img_filename)

                        with open(img_path, "wb") as f:
                            f.write(image_part.blob)

                        doc.images.append(ImageAsset(
                            path=img_path,
                            alt_text=f"插图-{img_count+1}",
                            source=f"第{para_idx+1}段",
                        ))
                        img_count += 1
                    except (KeyError, AttributeError):
                        continue

    # 也从页眉页脚提取图片
    for section in document.sections:
        for header in (section.header, section.footer):
            if header is None:
                continue
            for para in header.paragraphs:
                for run in para.runs:
                    drawing_elements = run._element.findall(qn("w:drawing"))
                    for drawing in drawing_elements:
                        blip = drawing.findall(".//" + qn("a:blip"))
                        for b in blip:
                            embed_id = b.get(qn("r:embed"))
                            if not embed_id:
                                continue
                            try:
                                image_part = document.part.related_parts[embed_id]
                                ext = _image_ext_from_content_type(image_part.content_type)
                                img_filename = f"image_{img_count:03d}.{ext}"
                                img_path = os.path.join(assets_dir, img_filename)
                                with open(img_path, "wb") as f:
                                    f.write(image_part.blob)
                                doc.images.append(ImageAsset(
                                    path=img_path,
                                    alt_text=f"页眉页脚图-{img_count+1}",
                                    source="页眉/页脚",
                                ))
                                img_count += 1
                            except (KeyError, AttributeError):
                                continue


def _image_ext_from_content_type(content_type: str) -> str:
    m = {
        "image/png": "png",
        "image/jpeg": "jpg",
        "image/gif": "gif",
        "image/bmp": "bmp",
        "image/tiff": "tiff",
        "image/svg+xml": "svg",
    }
    return m.get(content_type, "png")


# ── .pptx 读取 ────────────────────────────────────────────


def _read_pptx(file_path: str, doc: RawDocument, assets_dir: str, extract_images: bool = False):
    """用 python-pptx 读取 .pptx"""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation(file_path)
    texts = []
    img_count = 0

    for slide_idx, slide in enumerate(prs.slides):
        slide_label = f"--- slide {slide_idx + 1} ---"
        texts.append(slide_label)

        for shape in slide.shapes:
            # 文本
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)

            # 图片（可选，默认关闭）
            if extract_images and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    os.makedirs(assets_dir, exist_ok=True)
                    img = shape.image
                    ext = _image_ext_from_content_type(img.content_type)
                    img_filename = f"image_{img_count:03d}.{ext}"
                    img_path = os.path.join(assets_dir, img_filename)
                    with open(img_path, "wb") as f:
                        f.write(img.blob)
                    doc.images.append(ImageAsset(
                        path=img_path,
                        alt_text=shape.image_format.filename or f"幻灯片图片-{img_count+1}",
                        source=f"第{slide_idx+1}页",
                    ))
                    img_count += 1
                except Exception:
                    pass

            # 表格
            if shape.has_table:
                tbl = shape.table
                td = TableData()
                for row_idx, row in enumerate(tbl.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    if row_idx == 0:
                        td.headers = cells
                    else:
                        td.rows.append(cells)
                if td.headers or td.rows:
                    doc.tables.append(td)

        texts.append("")

    doc.text_content = "\n".join(texts)


# ── .xlsx 读取 ────────────────────────────────────────────


def _read_xlsx(file_path: str, doc: RawDocument):
    """用 openpyxl 读取 .xlsx（非 read_only 模式避免卡死）"""
    from openpyxl import load_workbook

    wb = load_workbook(file_path, data_only=True)
    parts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"## Sheet: {sheet_name}\n")
        table = TableData()
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            cells = [str(c) if c is not None else "" for c in row]
            line = " | ".join(cells)
            parts.append(line)
            if i == 0:
                table.headers = cells
            else:
                table.rows.append(cells)
        if table.headers:
            doc.tables.append(table)
        parts.append("")

    doc.text_content = "\n".join(parts)
    wb.close()


def _read_xls(file_path: str, doc: RawDocument):
    """用 xlrd 读取 .xls（降级方案）"""
    try:
        import xlrd
        wb = xlrd.open_workbook(file_path)
        parts = []
        for sheet_name in wb.sheet_names():
            ws = wb.sheet_by_name(sheet_name)
            parts.append(f"## Sheet: {sheet_name}\n")
            table = TableData()
            for i in range(ws.nrows):
                cells = [str(ws.cell_value(i, c)) for c in range(ws.ncols)]
                line = " | ".join(cells)
                parts.append(line)
                if i == 0:
                    table.headers = cells
                else:
                    table.rows.append(cells)
            if table.headers:
                doc.tables.append(table)
            parts.append("")
        doc.text_content = "\n".join(parts)
    except ImportError:
        doc.text_content = f"[无法读取]: .xls 需要 xlrd: pip install xlrd"


# ── PDF 读取（LibreOffice 转 .docx） ──────────────────────


def _read_pdf(file_path: str, doc: RawDocument, assets_dir: str, extract_images: bool = False):
    """PDF: LibreOffice → .docx → python-docx"""
    if not shutil.which("libreoffice"):
        doc.text_content = (
            "[无法读取]: PDF 需安装 LibreOffice 转换\n"
            "  sudo apt install libreoffice\n"
            "  提示: 可先手动将 PDF 转为 .docx 再传入"
        )
        return

    with tempfile.TemporaryDirectory() as tmpdir:
        try:
            result = subprocess.run(
                [
                    "libreoffice", "--headless", "--convert-to", "docx",
                    file_path, "--outdir", tmpdir,
                ],
                capture_output=True, text=True, timeout=120,
            )
            if result.returncode != 0:
                doc.text_content = (
                    f"[PDF转换失败]: LibreOffice 错误:\n{result.stderr.strip()}"
                )
                # 检查是否加密
                if "password" in result.stderr.lower() or "encrypt" in result.stderr.lower():
                    doc.text_content += (
                        "\n⚠️ 该 PDF 可能设置了打开密码，需要用户手动解除密码后重试"
                    )
                return

            tmp_files = list(Path(tmpdir).glob("*.docx"))
            if not tmp_files:
                doc.text_content = "[PDF转换失败]: 未生成 .docx 文件（可能是扫描件PDF）"
                return

            _read_docx(str(tmp_files[0]), doc, assets_dir, extract_images=extract_images)

        except subprocess.TimeoutExpired:
            doc.text_content = "[PDF转换失败]: LibreOffice 超时（>120秒）"
        except Exception as e:
            doc.text_content = f"[PDF转换异常]: {e}"


# ── 标题猜测 ──────────────────────────────────────────────


def _guess_title(text: str, file_path: str) -> str:
    """从文本中猜测文档标题"""
    if text:
        m = re.search(r"^#\s+(.+)", text, re.MULTILINE)
        if m:
            return m.group(1).strip()
        # 取第一段非空文本前 50 字
        for line in text.split("\n"):
            line = line.strip()
            if line and not line.startswith("#"):
                return line[:50]

    name = Path(file_path).stem
    name = re.sub(r"_[0-9]{8}", "", name)
    name = re.sub(r"\s+", " ", name).strip()
    return name


# ── 内容包中的图片路径转换（相对路径） ──────────────────────


def make_image_refs_relative(doc: RawDocument, content_pkg_dir: str) -> RawDocument:
    """将图片路径转为相对于内容包目录的路径"""
    if not doc.images:
        return doc

    try:
        base = os.path.abspath(content_pkg_dir)
        for img in doc.images:
            abs_path = os.path.abspath(img.path)
            try:
                rel = os.path.relpath(abs_path, base)
                img.path = rel
            except ValueError:
                pass  # 跨驱动器保留绝对路径
    except Exception:
        pass

    return doc